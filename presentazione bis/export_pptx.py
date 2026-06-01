#!/usr/bin/env python3
"""
Export presentazione_video_analysis.html → PPTX
Reply Template White style, 16:9, with embedded screenshots.
"""

import os
import re
from pathlib import Path
from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

# ── Reply Brand Colors ──
APPLE_GREEN = RGBColor(0x76, 0xB8, 0x2A)
DARK_GREEN = RGBColor(0x2D, 0x6A, 0x2E)
BRIGHT_GREEN = RGBColor(0x8C, 0xC6, 0x3F)
VIVID_BLUE = RGBColor(0x00, 0x72, 0xCE)
AMBER = RGBColor(0xF5, 0xA6, 0x23)
DEEP_ROSE = RGBColor(0xC4, 0x37, 0x6B)
BLACK = RGBColor(0x1A, 0x1A, 0x1A)
DARK_GRAY = RGBColor(0x55, 0x55, 0x55)
LIGHT_GRAY = RGBColor(0xF5, 0xF5, 0xF5)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# ── Slide dimensions 16:9 ──
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ── Paths ──
BASE_DIR = Path(__file__).parent
HTML_PATH = BASE_DIR / "presentazione_video_analysis.html"
SCREENSHOTS_DIR = BASE_DIR / "screenshots"
OUTPUT_PATH = BASE_DIR / "presentazione_video_analysis.pptx"


def clean(text):
    """Clean HTML text: strip whitespace, collapse newlines."""
    if not text:
        return ""
    text = re.sub(r'\s+', ' ', text.strip())
    return text


def add_textbox(slide, left, top, width, height):
    """Add a textbox and return its text_frame."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    return tf


def set_paragraph(tf, text, font_name="Arial", font_size=18, bold=False,
                  color=BLACK, alignment=PP_ALIGN.LEFT, space_after=Pt(6)):
    """Configure a paragraph in a text frame."""
    p = tf.paragraphs[0] if not tf.paragraphs[0].text else tf.add_paragraph()
    p.text = text
    p.font.name = font_name
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    p.space_after = space_after
    return p


def add_paragraph(tf, text, font_name="Arial", font_size=18, bold=False,
                  color=BLACK, alignment=PP_ALIGN.LEFT, space_after=Pt(6)):
    """Add a new paragraph to an existing text frame."""
    p = tf.add_paragraph()
    p.text = text
    p.font.name = font_name
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    p.space_after = space_after
    return p


def add_bullet(tf, text, font_size=16, color=BLACK, indent_level=0):
    """Add a bullet point."""
    p = tf.add_paragraph()
    p.text = text
    p.font.name = "Arial"
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.level = indent_level
    p.space_after = Pt(4)
    # Bullet character
    pPr = p._p.get_or_add_pPr()
    from pptx.oxml.ns import qn
    buNone = pPr.find(qn('a:buNone'))
    if buNone is not None:
        pPr.remove(buNone)
    buChar = pPr.makeelement(qn('a:buChar'), {'char': '▪'})
    buClr = pPr.makeelement(qn('a:buClr'), {})
    srgbClr = buClr.makeelement(qn('a:srgbClr'), {'val': '8CC63F'})
    buClr.append(srgbClr)
    buSzPct = pPr.makeelement(qn('a:buSzPct'), {'val': '130000'})
    pPr.append(buClr)
    pPr.append(buSzPct)
    pPr.append(buChar)
    return p


def set_slide_bg(slide, r, g, b):
    """Set solid background color for a slide."""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(r, g, b)


def set_slide_gradient_bg(slide, color1, color2):
    """Set a gradient background."""
    background = slide.background
    fill = background.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = color1
    fill.gradient_stops[0].position = 0.0
    fill.gradient_stops[1].color.rgb = color2
    fill.gradient_stops[1].position = 1.0


def add_image_fit(slide, img_path, left, top, max_width, max_height):
    """Add image preserving aspect ratio within bounds."""
    if not img_path.exists():
        return None
    with Image.open(img_path) as img:
        img_w, img_h = img.size
    aspect = img_w / img_h
    # Fit within max bounds
    if max_width / max_height > aspect:
        height = max_height
        width = int(height * aspect / Emu(1)) * Emu(1)
        width = min(Emu(int(max_height * aspect)), max_width)
        height = max_height
    else:
        width = max_width
        height = min(Emu(int(max_width / aspect)), max_height)
    pic = slide.shapes.add_picture(str(img_path), left, top, width, height)
    return pic


def make_cover(prs):
    """Slide 1: Cover with gradient background."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    set_slide_gradient_bg(slide, DARK_GREEN, APPLE_GREEN)

    # Badge
    tf = add_textbox(slide, Inches(1), Inches(3.5), Inches(6), Inches(0.5))
    set_paragraph(tf, "Copilot Studio + Work IQ MCP", "Arial", 12, False, WHITE, PP_ALIGN.LEFT)

    # Title
    tf = add_textbox(slide, Inches(1), Inches(4.0), Inches(10), Inches(2))
    set_paragraph(tf, "PROJECT REQUEST\nTRACKER", "Arial Black", 44, True, WHITE, PP_ALIGN.LEFT, Pt(0))

    # Subtitle
    tf = add_textbox(slide, Inches(1), Inches(6.0), Inches(10), Inches(0.8))
    set_paragraph(tf, "Analisi della demo video — Prompt, Tool e Flusso Operativo", "Arial", 16, False, WHITE)
    add_paragraph(tf, "Giugno 2026 · Video Analysis Report", "Arial", 13, False, RGBColor(0xCC, 0xCC, 0xCC))


def make_agenda(prs):
    """Slide 2: Agenda."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    tf = add_textbox(slide, Inches(1), Inches(0.6), Inches(11.3), Inches(0.7))
    set_paragraph(tf, "AGENDA", "Arial Black", 28, True, BLACK, PP_ALIGN.CENTER)
    add_paragraph(tf, "STRUTTURA DELLA PRESENTAZIONE", "Arial Black", 13, True, DARK_GRAY, PP_ALIGN.CENTER)

    items = [
        "Overview del Progetto",
        "System Prompt & Instructions",
        "Tool — Work IQ MCP",
        "Demo: Flusso Conversazionale",
        "Demo: AI Search & Risultati",
        "Demo: OnePager & Output Finale",
    ]
    tf = add_textbox(slide, Inches(3.5), Inches(2.0), Inches(6.5), Inches(5))
    for i, item in enumerate(items):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        run_num = p.add_run()
        run_num.text = f"{i+1}  "
        run_num.font.name = "Arial Black"
        run_num.font.size = Pt(22)
        run_num.font.color.rgb = APPLE_GREEN
        run_text = p.add_run()
        run_text.text = item
        run_text.font.name = "Arial"
        run_text.font.size = Pt(18)
        run_text.font.color.rgb = BLACK
        p.space_after = Pt(14)


def make_section(prs, icon, phrase, impact):
    """Section divider slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    tf = add_textbox(slide, Inches(2), Inches(2.2), Inches(9.3), Inches(0.8))
    set_paragraph(tf, icon, "Arial", 36, False, BLACK, PP_ALIGN.CENTER)

    tf = add_textbox(slide, Inches(2), Inches(3.2), Inches(9.3), Inches(0.8))
    set_paragraph(tf, phrase, "Arial Black", 32, True, APPLE_GREEN, PP_ALIGN.CENTER)

    tf = add_textbox(slide, Inches(2.5), Inches(4.2), Inches(8.3), Inches(1.5))
    set_paragraph(tf, impact, "Arial", 17, False, DARK_GRAY, PP_ALIGN.CENTER, Pt(8))


def make_flow_kpi(prs):
    """Slide 4: Flow diagram + KPI."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    tf = add_textbox(slide, Inches(0.5), Inches(0.4), Inches(12.3), Inches(0.7))
    set_paragraph(tf, "FLUSSO OPERATIVO END-TO-END", "Arial Black", 26, True, BLACK, PP_ALIGN.CENTER)
    add_paragraph(tf, "DALL'INPUT UTENTE ALL'OUTPUT DOCUMENTALE", "Arial Black", 12, True, DARK_GRAY, PP_ALIGN.CENTER)

    # Flow steps as boxes
    steps = [
        ("💬", "Prompt Utente", "Richiesta in linguaggio naturale"),
        ("📋", "Intake", "Raccolta 12 campi strutturati"),
        ("🔍", "AI Search", "Work IQ MCP copilot_chat"),
        ("🤔", "Decisione", "Procedi / Rivedi scope"),
        ("⚡", "Power Automate", "Crea request + OnePager"),
    ]
    box_w = Inches(2.1)
    box_h = Inches(1.3)
    arrow_w = Inches(0.35)
    total_w = len(steps) * box_w + (len(steps) - 1) * arrow_w
    start_x = (SLIDE_W - total_w) // 2
    y = Inches(2.0)

    for i, (icon, label, desc) in enumerate(steps):
        x = start_x + i * (box_w + arrow_w)
        # Box shape
        shape = slide.shapes.add_shape(
            1, x, y, box_w, box_h  # MSO_SHAPE.RECTANGLE
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = LIGHT_GRAY
        shape.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
        shape.line.width = Pt(1.5)

        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = icon
        p.font.size = Pt(20)
        p.alignment = PP_ALIGN.CENTER
        p2 = tf.add_paragraph()
        p2.text = label
        p2.font.name = "Arial Black"
        p2.font.size = Pt(10)
        p2.font.bold = True
        p2.alignment = PP_ALIGN.CENTER
        p3 = tf.add_paragraph()
        p3.text = desc
        p3.font.name = "Arial"
        p3.font.size = Pt(8)
        p3.font.color.rgb = DARK_GRAY
        p3.alignment = PP_ALIGN.CENTER

        # Arrow between steps
        if i < len(steps) - 1:
            ax = x + box_w
            atf = add_textbox(slide, ax, y + Inches(0.4), arrow_w, Inches(0.5))
            set_paragraph(atf, "→", "Arial", 22, True, APPLE_GREEN, PP_ALIGN.CENTER)

    # KPI cards
    kpis = [
        ("13:12", "Minuti di Demo"),
        ("12", "Campi Raccolti"),
        ("3", "Documenti Trovati"),
        ("55s", "AI Search Time"),
    ]
    kpi_w = Inches(2.5)
    kpi_h = Inches(1.3)
    kpi_gap = Inches(0.4)
    total_kpi = len(kpis) * kpi_w + (len(kpis) - 1) * kpi_gap
    kpi_x = (SLIDE_W - total_kpi) // 2
    kpi_y = Inches(4.5)

    for i, (value, label) in enumerate(kpis):
        x = kpi_x + i * (kpi_w + kpi_gap)
        shape = slide.shapes.add_shape(1, x, kpi_y, kpi_w, kpi_h)
        shape.fill.solid()
        shape.fill.fore_color.rgb = LIGHT_GRAY
        shape.line.color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
        shape.line.width = Pt(1)
        # Round corners
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = value
        p.font.name = "Arial Black"
        p.font.size = Pt(32)
        p.font.color.rgb = APPLE_GREEN
        p.alignment = PP_ALIGN.CENTER
        p.space_after = Pt(4)
        p2 = tf.add_paragraph()
        p2.text = label
        p2.font.name = "Arial"
        p2.font.size = Pt(11)
        p2.font.color.rgb = DARK_GRAY
        p2.alignment = PP_ALIGN.CENTER


def make_system_instructions(prs):
    """Slide 6: System Instructions prompt block."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    tf = add_textbox(slide, Inches(0.5), Inches(0.4), Inches(12.3), Inches(0.7))
    set_paragraph(tf, "SYSTEM INSTRUCTIONS", "Arial Black", 26, True, BLACK, PP_ALIGN.CENTER)
    add_paragraph(tf, "IL PROMPT DI SISTEMA DELL'AGENTE COPILOT STUDIO", "Arial Black", 12, True, DARK_GRAY, PP_ALIGN.CENTER)

    # Dark code block
    shape = slide.shapes.add_shape(1, Inches(1), Inches(1.6), Inches(11.3), Inches(5.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x1E, 0x1E, 0x2E)
    shape.line.color.rgb = APPLE_GREEN
    shape.line.width = Pt(3)

    tf = shape.text_frame
    tf.word_wrap = True

    lines = [
        ("Goal: ", True, "Collect all required request fields first, then run the"),
        ("", False, "AI/WorkIQ overlap check, then proceed based on user choice,"),
        ("", False, "and only then trigger the Power Automate creation flow."),
        ("", False, ""),
        ("// Step 1 — Intake", False, ""),
        ("Intake first: ", True, "DO NOT DIRECTLY ASK values to the user,"),
        ("", False, 'but MAKE SURE TO LAUNCH the "New project request" topic...'),
        ("", False, "it asks: Title, Requester, Sponsor, Objective/Outcome,"),
        ("", False, "Benefits, Urgency, Deadline, Impacts, Dependencies, Effort, Risks, Notes"),
        ("", False, ""),
        ("// Step 2 — AI Search", False, ""),
        ("Then run AI search ", True, "(Work IQ copilot_chat MCP tool)"),
        ("", False, "to detect overlap and dependencies."),
        ("", False, ""),
        ("// Step 3 — Decisione utente", False, ""),
        ('Present findings: ', True, '"What would you like to do next?"'),
        ("", False, "  ✅ Proceed as-is"),
        ("", False, "  ⟲ Revise the scope"),
        ("", False, ""),
        ("// Step 4 — Creazione", False, ""),
        ("Only after confirmation ", True, "→ trigger Power Automate"),
        ("", False, '"Create new project request" flow.'),
    ]

    for i, (keyword, is_keyword, rest) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if keyword:
            run_k = p.add_run()
            run_k.text = keyword
            run_k.font.name = "Consolas"
            run_k.font.size = Pt(12)
            run_k.font.color.rgb = RGBColor(0xCB, 0xA6, 0xF7) if is_keyword else RGBColor(0x6C, 0x70, 0x86)
            run_k.font.bold = is_keyword
        if rest:
            run_r = p.add_run()
            run_r.text = rest
            run_r.font.name = "Consolas"
            run_r.font.size = Pt(12)
            run_r.font.color.rgb = RGBColor(0xCD, 0xD6, 0xF4)
        if not keyword and not rest:
            p.text = ""
            p.font.size = Pt(8)
        elif keyword.startswith("//"):
            p.runs[0].font.color.rgb = RGBColor(0x6C, 0x70, 0x86)
            p.runs[0].font.italic = True
        p.space_after = Pt(2)


def make_image_side(prs, img_filename, evidence, title, body_lines, bullets=None, footer=None):
    """Image + Side Text slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Image on left half
    img_path = SCREENSHOTS_DIR / img_filename
    add_image_fit(slide, img_path, Inches(0.3), Inches(0.3), Inches(6.5), Inches(6.9))

    # Text on right
    x_text = Inches(7.1)
    w_text = Inches(5.8)

    y = Inches(1.5)
    # Evidence label
    tf = add_textbox(slide, x_text, y, w_text, Inches(0.4))
    set_paragraph(tf, evidence, "Arial Black", 11, True, BRIGHT_GREEN, PP_ALIGN.LEFT)
    y += Inches(0.5)

    # Title
    tf = add_textbox(slide, x_text, y, w_text, Inches(0.6))
    set_paragraph(tf, title, "Arial Black", 20, True, BLACK, PP_ALIGN.LEFT)
    y += Inches(0.7)

    # Body lines
    if body_lines:
        tf = add_textbox(slide, x_text, y, w_text, Inches(1.0))
        for i, line in enumerate(body_lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            p.font.name = "Arial"
            p.font.size = Pt(14)
            p.font.color.rgb = BLACK
            p.space_after = Pt(6)
        y += Inches(0.3 + 0.25 * len(body_lines))

    # Bullets
    if bullets:
        tf = add_textbox(slide, x_text, y, w_text, Inches(3.0))
        for i, bullet_text in enumerate(bullets):
            if i == 0:
                add_bullet(tf, bullet_text, 13, BLACK)
                tf.paragraphs[0].text = ""  # clear default empty
            else:
                add_bullet(tf, bullet_text, 13, BLACK)
        y += Inches(0.3 + 0.22 * len(bullets))

    # Footer note
    if footer:
        tf = add_textbox(slide, x_text, y + Inches(0.2), w_text, Inches(0.8))
        shape = tf._txBody.getparent()
        # Add background
        set_paragraph(tf, footer, "Arial", 11, False, DARK_GRAY, PP_ALIGN.LEFT)


def make_user_prompt(prs):
    """Slide 8: User prompt."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    tf = add_textbox(slide, Inches(0.5), Inches(0.4), Inches(12.3), Inches(0.7))
    set_paragraph(tf, "PROMPT UTENTE", "Arial Black", 26, True, BLACK, PP_ALIGN.CENTER)
    add_paragraph(tf, "IL PROMPT CHE AVVIA IL PROCESSO", "Arial Black", 12, True, DARK_GRAY, PP_ALIGN.CENTER)

    # Prompt quote box
    shape = slide.shapes.add_shape(1, Inches(2), Inches(1.8), Inches(9.3), Inches(1.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0x1E, 0x1E, 0x2E)
    shape.line.color.rgb = VIVID_BLUE
    shape.line.width = Pt(3)

    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = '"I want to open a new internal project request to improve\ncustomer data quality in CRM and reduce duplicate\ncustomer records"'
    p.font.name = "Consolas"
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(0xA6, 0xE3, 0xA1)
    p.alignment = PP_ALIGN.CENTER

    # Screenshot
    img_path = SCREENSHOTS_DIR / "03_prompt_user_request.png"
    add_image_fit(slide, img_path, Inches(1.5), Inches(3.6), Inches(10.3), Inches(3.5))

    # Caption
    tf = add_textbox(slide, Inches(2), Inches(7.0), Inches(9.3), Inches(0.4))
    set_paragraph(tf, 'L\'agente interpreta il prompt e instrada al topic "New project request"',
                  "Arial", 11, False, DARK_GRAY, PP_ALIGN.CENTER)


def make_tool_table(prs):
    """Slide 11: Tool map table."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    tf = add_textbox(slide, Inches(0.5), Inches(0.4), Inches(12.3), Inches(0.7))
    set_paragraph(tf, "MAPPA DEI TOOL", "Arial Black", 26, True, BLACK, PP_ALIGN.CENTER)
    add_paragraph(tf, "TUTTI I COMPONENTI COINVOLTI NEL FLUSSO", "Arial Black", 12, True, DARK_GRAY, PP_ALIGN.CENTER)

    # Table
    rows = 6
    cols = 4
    table_shape = slide.shapes.add_table(rows, cols, Inches(1), Inches(1.8), Inches(11.3), Inches(4.5))
    table = table_shape.table

    # Column widths
    table.columns[0].width = Inches(3.5)
    table.columns[1].width = Inches(2.0)
    table.columns[2].width = Inches(4.3)
    table.columns[3].width = Inches(1.5)

    # Header row
    headers = ["Tool", "Tipo", "Funzione", "Durata"]
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = APPLE_GREEN
        p = cell.text_frame.paragraphs[0]
        p.font.name = "Arial Black"
        p.font.size = Pt(12)
        p.font.color.rgb = WHITE
        p.font.bold = True

    # Data rows
    data = [
        ("copilot_chat", "MCP", "AI Search nella KnowledgeBase OneDrive", "55.31s"),
        ("New project request", "Topic", "Raccolta 12 campi strutturati (intake form)", "~3 min"),
        ("Revise Project Request Title", "Topic", "Revisione titolo su richiesta utente", "~15s"),
        ("Create new project request", "Power Automate", "Genera OnePager su SharePoint + Request ID", "10.11s"),
        ("Copilot in Word", "M365 Copilot", "Formattazione AI Notes da markdown a Word", "~20s"),
    ]
    type_colors = {
        "MCP": VIVID_BLUE,
        "Topic": APPLE_GREEN,
        "Power Automate": AMBER,
        "M365 Copilot": AMBER,
    }
    for r, (tool, tipo, func, dur) in enumerate(data, 1):
        for c, val in enumerate([tool, tipo, func, dur]):
            cell = table.cell(r, c)
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            p.font.name = "Arial"
            p.font.size = Pt(13)
            p.font.color.rgb = BLACK
            if c == 0:
                p.font.bold = True
            if c == 1:
                p.font.color.rgb = type_colors.get(tipo, BLACK)
                p.font.bold = True
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY


def make_dual_screenshot(prs, title, subtitle, img1, caption1, img2, caption2):
    """Two screenshots side by side."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    tf = add_textbox(slide, Inches(0.5), Inches(0.4), Inches(12.3), Inches(0.7))
    set_paragraph(tf, title, "Arial Black", 26, True, BLACK, PP_ALIGN.CENTER)
    add_paragraph(tf, subtitle, "Arial Black", 12, True, DARK_GRAY, PP_ALIGN.CENTER)

    # Left image
    add_image_fit(slide, SCREENSHOTS_DIR / img1, Inches(0.5), Inches(1.8), Inches(6.0), Inches(4.5))
    tf = add_textbox(slide, Inches(0.5), Inches(6.5), Inches(6.0), Inches(0.4))
    set_paragraph(tf, caption1, "Arial", 11, False, DARK_GRAY, PP_ALIGN.CENTER)

    # Right image
    add_image_fit(slide, SCREENSHOTS_DIR / img2, Inches(6.8), Inches(1.8), Inches(6.0), Inches(4.5))
    tf = add_textbox(slide, Inches(6.8), Inches(6.5), Inches(6.0), Inches(0.4))
    set_paragraph(tf, caption2, "Arial", 11, False, DARK_GRAY, PP_ALIGN.CENTER)


def make_video_slide(prs, title, subtitle, screenshot_file, clip_name, duration):
    """Video phase slide — shows screenshot poster + info since PPTX can't embed web video."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    tf = add_textbox(slide, Inches(0.5), Inches(0.4), Inches(12.3), Inches(0.7))
    set_paragraph(tf, title, "Arial Black", 24, True, BLACK, PP_ALIGN.CENTER)
    add_paragraph(tf, subtitle, "Arial Black", 11, True, DARK_GRAY, PP_ALIGN.CENTER)

    # Video label bar
    shape = slide.shapes.add_shape(1, Inches(1.5), Inches(1.7), Inches(10.3), Inches(0.45))
    shape.fill.solid()
    shape.fill.fore_color.rgb = APPLE_GREEN
    shape.line.fill.background()
    tf = shape.text_frame
    p = tf.paragraphs[0]
    run1 = p.add_run()
    run1.text = f"▶  {clip_name}"
    run1.font.name = "Arial Black"
    run1.font.size = Pt(11)
    run1.font.color.rgb = WHITE
    run2 = p.add_run()
    run2.text = f"                                                                         {duration}"
    run2.font.name = "Arial"
    run2.font.size = Pt(11)
    run2.font.color.rgb = WHITE

    # Screenshot as poster
    img_path = SCREENSHOTS_DIR / screenshot_file
    add_image_fit(slide, img_path, Inches(1.5), Inches(2.15), Inches(10.3), Inches(4.8))

    # Note
    tf = add_textbox(slide, Inches(2), Inches(7.0), Inches(9.3), Inches(0.4))
    set_paragraph(tf, f"📎 Video clip disponibile: clips/{clip_name.lower().replace(' ', '_')}.mp4",
                  "Arial", 10, False, DARK_GRAY, PP_ALIGN.CENTER)


def make_ai_results(prs):
    """AI Search results detail slide."""
    make_image_side(prs,
        "10_demo_ai_search_results.png",
        "TOP RELATED DOCUMENTS",
        "RISULTATI AI SEARCH",
        [],
        [
            "Customer_Master_Cleanup_Initiative.docx — Strategic initiative, duplicati, naming, VAT IDs",
            "CRM_Data_Quality_Meeting_Notes.docx — Stakeholder discussions, validation rules, duplicate detection",
            "Customer_Data_Risk_Assessment.docx — Risk analysis, CRM/ERP sync failures, ownership gaps",
        ],
        "⚠️ Temi emergenti: governance già discussa, rischi ERP sync, stakeholder Sales Ops/IT/Data Governance"
    )


def make_before_after(prs):
    """OnePager before/after."""
    make_dual_screenshot(prs,
        "ONEPAGER — PRIMA E DOPO",
        "DA MARKDOWN RAW A DOCUMENTO FORMATTATO CON COPILOT IN WORD",
        "12_demo_onepager_raw.png", "AI Notes in formato markdown grezzo",
        "14_demo_final_onepager.png", "Documento finale formattato da Copilot Word"
    )


def make_statement(prs):
    """Statement slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    tf = add_textbox(slide, Inches(1.5), Inches(2.0), Inches(10.3), Inches(2.5))
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    p = tf.paragraphs[0]
    run1 = p.add_run()
    run1.text = "DALL'IDEA AL\n"
    run1.font.name = "Arial Black"
    run1.font.size = Pt(40)
    run1.font.color.rgb = BLACK

    p2 = tf.add_paragraph()
    run2 = p2.add_run()
    run2.text = "DOCUMENTO UFFICIALE"
    run2.font.name = "Arial Black"
    run2.font.size = Pt(40)
    run2.font.color.rgb = APPLE_GREEN
    p2.alignment = PP_ALIGN.CENTER

    p3 = tf.add_paragraph()
    run3 = p3.add_run()
    run3.text = "IN 13 MINUTI"
    run3.font.name = "Arial Black"
    run3.font.size = Pt(40)
    run3.font.color.rgb = BLACK
    p3.alignment = PP_ALIGN.CENTER

    tf2 = add_textbox(slide, Inches(2.5), Inches(5.0), Inches(8.3), Inches(1.2))
    set_paragraph(tf2,
        "Un singolo prompt attiva intake strutturato, AI search nella knowledge base aziendale, "
        "e generazione automatica su SharePoint — senza uscire dalla chat.",
        "Arial", 15, False, DARK_GRAY, PP_ALIGN.CENTER
    )


def make_thankyou(prs):
    """Thank you slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_gradient_bg(slide, DARK_GREEN, APPLE_GREEN)

    tf = add_textbox(slide, Inches(1), Inches(2.5), Inches(11.3), Inches(1.5))
    set_paragraph(tf, "THANK YOU", "Arial Black", 60, True, WHITE, PP_ALIGN.CENTER)

    tf = add_textbox(slide, Inches(2), Inches(4.5), Inches(9.3), Inches(0.8))
    set_paragraph(tf, "Project Request Tracker — Copilot Studio + Work IQ MCP",
                  "Arial", 16, False, RGBColor(0xDD, 0xDD, 0xDD), PP_ALIGN.CENTER)


def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # 1. Cover
    make_cover(prs)

    # 2. Agenda
    make_agenda(prs)

    # 3. Section: Overview
    make_section(prs, "🎯", "OVERVIEW",
        "Un agente Copilot Studio che gestisce l'intero ciclo di creazione project request — "
        "dall'intake conversazionale alla generazione automatica di un OnePager su SharePoint.")

    # 4. Flow + KPI
    make_flow_kpi(prs)

    # 5. Section: Prompt
    make_section(prs, "📝", "SYSTEM PROMPT",
        "Le istruzioni di sistema definiscono il comportamento dell'agente: "
        "raccolta dati, AI search, decisione utente e trigger Power Automate.")

    # 6. System Instructions
    make_system_instructions(prs)

    # 7. Instructions screenshot
    make_image_side(prs,
        "01_prompt_system_instructions.png",
        "COPILOT STUDIO — OVERVIEW",
        "INSTRUCTIONS DELL'AGENTE",
        ["Le System Instructions definiscono un workflow in 4 fasi rigorosamente sequenziali:"],
        [
            "Intake: raccolta di 12 campi strutturati via topic dedicato",
            "AI Search: ricerca nella KnowledgeBase OneDrive via Work IQ MCP",
            "Decisione: presentazione risultati e scelta utente",
            "Creazione: Power Automate genera OnePager su SharePoint",
        ],
        "Modello: GPT-5 Chat"
    )

    # 8. User prompt
    make_user_prompt(prs)

    # 9. Section: Tool
    make_section(prs, "🔧", "TOOL UTILIZZATI",
        "Work IQ MCP, Topics Copilot Studio, AI Tool Node e Power Automate orchestrano il flusso end-to-end.")

    # 10. Work IQ MCP
    make_image_side(prs,
        "04_tool_workiq_mcp_panel.png",
        "MODEL CONTEXT PROTOCOL",
        "WORK IQ COPILOT (PREVIEW)",
        ["Server MCP che espone il tool copilot_chat per interrogare Microsoft 365:"],
        [
            "Ricerca semantica su documenti, email, chat, siti",
            "Accesso a OneDrive, SharePoint, Teams, Mail",
            "Priorità ai tool workload-specific quando il contesto è chiaro",
        ],
        "Query: \"Search inside my OneDrive folder 'KnowledgeBase' for documents related to CRM data quality...\""
    )

    # 11. Tool table
    make_tool_table(prs)

    # 12. copilot_chat before/after
    make_dual_screenshot(prs,
        "COPILOT_CHAT — ESECUZIONE",
        "55 SECONDI DI AI SEARCH NELLA KNOWLEDGEBASE",
        "05_tool_copilot_chat_working.png", "Tool in esecuzione — stato \"Working\"",
        "06_tool_copilot_chat_completed.png", "Tool completato — risultati disponibili"
    )

    # 13. Section: Demo
    make_section(prs, "🎬", "DEMO VIDEO",
        "Registrazioni del flusso operativo completo: dal prompt iniziale alla generazione del documento finale.")

    # 14. Video: Setup
    make_video_slide(prs, "FASE 1 — SETUP AGENTE",
        "CONFIGURAZIONE COPILOT STUDIO: OVERVIEW, INSTRUCTIONS, TOOLS",
        "01_prompt_system_instructions.png", "Setup Agente & Instructions", "3:40")

    # 15. Video: Prompt + Data Collection
    make_video_slide(prs, "FASE 2 — PROMPT & RACCOLTA DATI",
        "L'UTENTE INVIA IL PROMPT E L'AGENTE RACCOGLIE I 12 CAMPI",
        "03_prompt_user_request.png", "Prompt Utente & Data Collection", "2:40")

    # 16. Video: AI Search
    make_video_slide(prs, "FASE 3 — AI SEARCH CON WORK IQ",
        "IL TOOL COPILOT_CHAT CERCA NELLA KNOWLEDGEBASE E RESTITUISCE INSIGHT",
        "10_demo_ai_search_results.png", "AI Search & Risultati KnowledgeBase", "2:00")

    # 17. AI Results detail
    make_ai_results(prs)

    # 18. Video: Revise & Create
    make_video_slide(prs, "FASE 4 — REVISIONE & CREAZIONE",
        "L'UTENTE RIVEDE IL TITOLO, CONFERMA, E POWER AUTOMATE GENERA LA REQUEST",
        "09_tool_create_request_completed.png", "Revisione Titolo & Creazione Request", "1:10")

    # 19. Video: OnePager
    make_video_slide(prs, "FASE 5 — ONEPAGER SU SHAREPOINT",
        "DOCUMENTO GENERATO AUTOMATICAMENTE + FORMATTAZIONE CON COPILOT IN WORD",
        "14_demo_final_onepager.png", "OnePager & Copilot Word Formatting", "2:43")

    # 20. Before/After
    make_before_after(prs)

    # 21. Statement
    make_statement(prs)

    # 22. Knowledge Base
    make_image_side(prs,
        "15_demo_knowledge_source_doc.png",
        "KNOWLEDGE BASE",
        "DOCUMENTI SORGENTE",
        ["I documenti nella cartella OneDrive KnowledgeBase costituiscono il corpus per la ricerca semantica:"],
        [
            "Customer Master Cleanup Initiative — Piano strategico pulizia dati CRM",
            "CRM Data Quality Meeting Notes — Verbali riunioni stakeholder",
            "Customer Data Risk Assessment — Analisi rischi e mitigazioni",
        ],
        "L'agente cerca solo nella cartella specificata, garantendo risultati pertinenti e controllati."
    )

    # 23. Topics Editor
    make_image_side(prs,
        "07_tool_topics_editor.png",
        "COPILOT STUDIO — AUTHORING",
        "TOPICS EDITOR",
        ['Il topic "New project request" definisce il flusso di raccolta dati con nodi Question:'],
        [
            "Urgency → Global.varUrgency",
            "Expected Effort → Global.varEffort",
            "Deadline → Global.varDeadline",
        ],
        "Ogni variabile globale viene passata al flusso Power Automate per la generazione del OnePager."
    )

    # 24. Thank You
    make_thankyou(prs)

    prs.save(str(OUTPUT_PATH))
    print(f"✅ PPTX generato: {OUTPUT_PATH}")
    print(f"   Slide: {len(prs.slides)}")
    print(f"   Size: {OUTPUT_PATH.stat().st_size / 1024 / 1024:.1f} MB")


if __name__ == "__main__":
    main()
