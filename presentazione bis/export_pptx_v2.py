#!/usr/bin/env python3
"""
Export presentazione_video_analysis.html → PPTX  (v2)
Reply Template White — 16:9 — improved layout & verbatim prompt from video.
"""

import re
from pathlib import Path
from lxml import etree
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.opc.package import Part
from pptx.opc.packuri import PackURI
from PIL import Image

# ── Reply Brand Colors ──────────────────────────────────────────────
APPLE_GREEN  = RGBColor(0x76, 0xB8, 0x2A)
DARK_GREEN   = RGBColor(0x2D, 0x6A, 0x2E)
BRIGHT_GREEN = RGBColor(0x8C, 0xC6, 0x3F)
VIVID_BLUE   = RGBColor(0x00, 0x72, 0xCE)
AMBER        = RGBColor(0xF5, 0xA6, 0x23)
BLACK        = RGBColor(0x1A, 0x1A, 0x1A)
DARK_GRAY    = RGBColor(0x55, 0x55, 0x55)
MID_GRAY     = RGBColor(0x88, 0x88, 0x88)
LIGHT_GRAY   = RGBColor(0xF5, 0xF5, 0xF5)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
CODE_BG      = RGBColor(0x1E, 0x1E, 0x2E)
CODE_TEXT    = RGBColor(0xCD, 0xD6, 0xF4)
CODE_KW      = RGBColor(0xCB, 0xA6, 0xF7)
CODE_STR     = RGBColor(0xA6, 0xE3, 0xA1)
BORDER_GRAY  = RGBColor(0xDD, 0xDD, 0xDD)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

BASE_DIR       = Path(__file__).parent
SCREENSHOTS    = BASE_DIR / "screenshots"
CLIPS          = BASE_DIR / "clips"
OUTPUT         = BASE_DIR / "presentazione_video_analysis.pptx"

_video_counter = 0  # unique index for embedded video parts

# ── Helpers ──────────────────────────────────────────────────────────

def tb(slide, l, t, w, h, anchor="t"):
    """Create textbox, return text_frame. anchor: t/m/b."""
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    if anchor == "m":
        from pptx.enum.text import MSO_ANCHOR
        tf.paragraphs[0].alignment = PP_ALIGN.LEFT
        box.text_frame.auto_size = None
    return tf


def run(p, text, name="Arial", sz=18, bold=False, color=BLACK, italic=False):
    """Add a run to a paragraph."""
    r = p.add_run()
    r.text = text
    r.font.name = name
    r.font.size = Pt(sz)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return r


def para(tf, text, name="Arial", sz=18, bold=False, color=BLACK,
         align=PP_ALIGN.LEFT, after=Pt(6), first=False):
    """Add paragraph (or use first existing one)."""
    if first and not tf.paragraphs[0].text:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = text
    p.font.name = name
    p.font.size = Pt(sz)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    p.space_after = after
    return p


def bullet(tf, text, sz=14, color=BLACK, level=0):
    """Add green-square bullet."""
    p = tf.add_paragraph()
    p.text = text
    p.font.name = "Arial"
    p.font.size = Pt(sz)
    p.font.color.rgb = color
    p.level = level
    p.space_after = Pt(8)
    p.space_before = Pt(2)
    pPr = p._p.get_or_add_pPr()
    buChar = pPr.makeelement(qn('a:buChar'), {'char': '▪'})
    buClr  = pPr.makeelement(qn('a:buClr'), {})
    buClr.append(buClr.makeelement(qn('a:srgbClr'), {'val': '8CC63F'}))
    buSz   = pPr.makeelement(qn('a:buSzPct'), {'val': '130000'})
    pPr.append(buClr); pPr.append(buSz); pPr.append(buChar)
    return p


def img_fit(slide, path, l, t, max_w, max_h, border=False):
    """Add image preserving aspect ratio, centered in bounds."""
    if not path.exists():
        return None
    with Image.open(path) as im:
        iw, ih = im.size
    ratio = iw / ih
    if max_w / max_h > ratio:
        h = max_h; w = Emu(int(h * ratio))
    else:
        w = max_w; h = Emu(int(w / ratio))
    # center both horizontally and vertically
    offset_x = (max_w - w) // 2
    offset_y = (max_h - h) // 2
    pic = slide.shapes.add_picture(str(path), l + offset_x, t + offset_y, w, h)
    if border:
        pic.line.color.rgb = BORDER_GRAY
        pic.line.width = Pt(0.75)
    return pic


def embed_video(slide, prs, video_path, poster_path, left, top, max_w, max_h):
    """Embed an MP4 video with poster frame into a slide shape."""
    global _video_counter
    # Add poster image (used as thumbnail/poster frame)
    pic_shape = img_fit(slide, poster_path, left, top, max_w, max_h)
    if pic_shape is None or not video_path.exists():
        return pic_shape

    # Read video
    with open(str(video_path), 'rb') as f:
        video_blob = f.read()

    _video_counter += 1
    vname = PackURI(f'/ppt/media/video{_video_counter}.mp4')
    video_part = Part(vname, 'video/mp4', prs.part.package, video_blob)

    # Two relationships: video (for a:videoFile r:link) and media (for p14:media r:embed)
    RT_VIDEO = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships/video'
    RT_MEDIA = 'http://schemas.microsoft.com/office/2007/relationships/media'
    rId_video = slide.part.relate_to(video_part, RT_VIDEO)
    rId_media = slide.part.relate_to(video_part, RT_MEDIA)

    # Modify the pic element XML
    pic_el = pic_shape._element
    cNvPr = pic_el.find(qn('p:nvPicPr')).find(qn('p:cNvPr'))
    nvPr   = pic_el.find(qn('p:nvPicPr')).find(qn('p:nvPr'))

    # Add hlinkClick action so PowerPoint treats it as playable media
    hlinkClick = etree.SubElement(cNvPr, qn('a:hlinkClick'))
    hlinkClick.set(qn('r:id'), '')
    hlinkClick.set('action', 'ppaction://media')

    # a:videoFile references the video relationship
    videoFile = etree.SubElement(nvPr, qn('a:videoFile'))
    videoFile.set(qn('r:link'), rId_video)

    # p14:media extension for PowerPoint 2010+
    extLst = etree.SubElement(nvPr, qn('p:extLst'))
    ext = etree.SubElement(extLst, qn('p:ext'))
    ext.set('uri', '{DAA4B4D4-6D71-4841-9C94-3DE7FCFB9230}')
    p14_ns = 'http://schemas.microsoft.com/office/powerpoint/2010/main'
    media_el = etree.SubElement(ext, f'{{{p14_ns}}}media')
    media_el.set(qn('r:embed'), rId_media)

    return pic_shape


def bg_solid(slide, r, g, b):
    fill = slide.background.fill; fill.solid()
    fill.fore_color.rgb = RGBColor(r, g, b)


def bg_gradient(slide, c1, c2):
    fill = slide.background.fill; fill.gradient()
    fill.gradient_stops[0].color.rgb = c1; fill.gradient_stops[0].position = 0.0
    fill.gradient_stops[1].color.rgb = c2; fill.gradient_stops[1].position = 1.0


def title_bar(slide, title, subtitle=None):
    """Standard slide title + subtitle at top."""
    tf = tb(slide, Inches(0.8), Inches(0.45), Inches(11.7), Inches(0.65))
    para(tf, title, "Arial Black", 24, True, BLACK, PP_ALIGN.CENTER, Pt(2), first=True)
    if subtitle:
        para(tf, subtitle, "Arial Black", 11, True, DARK_GRAY, PP_ALIGN.CENTER, Pt(0))


def green_line(slide, y=None):
    """Thin green accent line across slide."""
    if y is None:
        y = Inches(1.25)
    shape = slide.shapes.add_shape(1, Inches(0.8), y, Inches(11.7), Pt(3))
    shape.fill.solid(); shape.fill.fore_color.rgb = APPLE_GREEN
    shape.line.fill.background()


# ── Slide builders ───────────────────────────────────────────────────

def slide_cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg_gradient(s, DARK_GREEN, APPLE_GREEN)

    # Badge pill
    tf = tb(s, Inches(1.4), Inches(3.4), Inches(5), Inches(0.4))
    para(tf, "COPILOT STUDIO  ·  WORK IQ MCP", "Arial", 12, True,
         RGBColor(0xDD,0xFF,0xDD), PP_ALIGN.LEFT, Pt(0), first=True)

    # Main title
    tf = tb(s, Inches(1.4), Inches(3.9), Inches(10), Inches(1.8))
    para(tf, "PROJECT REQUEST\nTRACKER", "Arial Black", 48, True, WHITE,
         PP_ALIGN.LEFT, Pt(6), first=True)

    # Subtitle
    tf = tb(s, Inches(1.4), Inches(5.9), Inches(10), Inches(0.5))
    para(tf, "Analisi della demo video — Prompt, Tool e Flusso Operativo",
         "Arial", 16, False, RGBColor(0xEE,0xEE,0xEE), PP_ALIGN.LEFT, Pt(0), first=True)

    tf = tb(s, Inches(1.4), Inches(6.5), Inches(10), Inches(0.4))
    para(tf, "Giugno 2026  ·  Video Analysis Report", "Arial", 13, False,
         RGBColor(0xBB,0xBB,0xBB), PP_ALIGN.LEFT, Pt(0), first=True)


def slide_agenda(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(s, "AGENDA", "STRUTTURA DELLA PRESENTAZIONE")
    green_line(s)

    items = [
        "Overview del Progetto",
        "System Prompt & Instructions",
        "Tool — Work IQ MCP",
        "Demo: Flusso Conversazionale",
        "Demo: AI Search & Risultati",
        "Demo: OnePager & Output Finale",
    ]
    tf = tb(s, Inches(3.4), Inches(2.0), Inches(6.5), Inches(5.1))
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run(p, f"{i+1}", "Arial Black", 28, True, APPLE_GREEN)
        run(p, f"   {item}", "Arial", 18, False, BLACK)
        p.space_after = Pt(28)
        p.space_before = Pt(4)
    # Decorative numbers column line
    shape = s.shapes.add_shape(1, Inches(4.2), Inches(1.9), Pt(2), Inches(5.0))
    shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(0xE0,0xE0,0xE0)
    shape.line.fill.background()


def slide_section(prs, icon, phrase, impact):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    tf = tb(s, Inches(2), Inches(2.3), Inches(9.3), Inches(0.8))
    para(tf, icon, "Arial", 44, False, BLACK, PP_ALIGN.CENTER, Pt(8), first=True)

    tf = tb(s, Inches(2), Inches(3.3), Inches(9.3), Inches(0.9))
    para(tf, phrase, "Arial Black", 34, True, APPLE_GREEN, PP_ALIGN.CENTER, Pt(8), first=True)

    # Decorative underline
    ul = s.shapes.add_shape(1, Inches(5.2), Inches(4.25), Inches(2.9), Pt(4))
    ul.fill.solid(); ul.fill.fore_color.rgb = BRIGHT_GREEN
    ul.line.fill.background()

    tf = tb(s, Inches(2.8), Inches(4.7), Inches(7.7), Inches(1.8))
    para(tf, impact, "Arial", 16, False, DARK_GRAY, PP_ALIGN.CENTER, Pt(0), first=True)


def slide_flow_kpi(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(s, "FLUSSO OPERATIVO END-TO-END",
              "DALL'INPUT UTENTE ALL'OUTPUT DOCUMENTALE")
    green_line(s)

    steps = [
        ("💬", "PROMPT UTENTE", "Richiesta in\nlinguaggio naturale"),
        ("📋", "INTAKE", "Raccolta 12 campi\nstrutturati"),
        ("🔍", "AI SEARCH", "Work IQ MCP\ncopilot_chat"),
        ("🤔", "DECISIONE", "Procedi /\nRivedi scope"),
        ("⚡", "POWER AUTOMATE", "Crea request\n+ OnePager"),
    ]
    box_w = Inches(2.1);  box_h = Inches(1.7)
    gap = Inches(0.45)
    total = len(steps) * box_w + (len(steps)-1) * gap
    x0 = (SLIDE_W - total) // 2
    y = Inches(1.8)

    for i, (icon, label, desc) in enumerate(steps):
        x = x0 + i * (box_w + gap)
        shape = s.shapes.add_shape(1, x, y, box_w, box_h)
        shape.fill.solid(); shape.fill.fore_color.rgb = LIGHT_GRAY
        shape.line.color.rgb = RGBColor(0xDE,0xDE,0xDE); shape.line.width = Pt(1)
        tf = shape.text_frame; tf.word_wrap = True
        tf.margin_top = Inches(0.12)
        para(tf, icon, "Arial", 26, False, BLACK, PP_ALIGN.CENTER, Pt(4), first=True)
        para(tf, label, "Arial Black", 9, True, BLACK, PP_ALIGN.CENTER, Pt(4))
        para(tf, desc, "Arial", 8, False, DARK_GRAY, PP_ALIGN.CENTER, Pt(0))

        if i < len(steps)-1:
            atf = tb(s, x + box_w, y + Inches(0.6), gap, Inches(0.4))
            para(atf, "→", "Arial", 26, True, APPLE_GREEN, PP_ALIGN.CENTER, Pt(0), first=True)

    # KPI row
    kpis = [("13:12", "Minuti di Demo"), ("12", "Campi Raccolti"),
            ("3", "Documenti Trovati"), ("55s", "AI Search Time")]
    kw = Inches(2.6); kh = Inches(1.5); kg = Inches(0.5)
    ktotal = len(kpis)*kw + (len(kpis)-1)*kg
    kx0 = (SLIDE_W - ktotal) // 2
    ky = Inches(4.8)

    for i, (val, lbl) in enumerate(kpis):
        x = kx0 + i*(kw+kg)
        shape = s.shapes.add_shape(1, x, ky, kw, kh)
        shape.fill.solid(); shape.fill.fore_color.rgb = WHITE
        shape.line.color.rgb = RGBColor(0xE0,0xE0,0xE0); shape.line.width = Pt(1.5)
        shape.shadow.inherit = False
        tf = shape.text_frame; tf.word_wrap = True
        tf.margin_top = Inches(0.1)
        para(tf, val, "Arial Black", 38, True, APPLE_GREEN, PP_ALIGN.CENTER, Pt(4), first=True)
        para(tf, lbl, "Arial", 12, False, DARK_GRAY, PP_ALIGN.CENTER, Pt(0))


def slide_prompt_verbatim(prs):
    """System Instructions — VERBATIM from video (2 slides for readability)."""

    # ── Slide A: Prompt Part 1 (Goal + Intake + AI Search) ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(s, "SYSTEM INSTRUCTIONS (1/2)", "PROMPT VERBATIM — COPILOT STUDIO · GPT-5 CHAT")
    green_line(s)

    shape = s.shapes.add_shape(1, Inches(0.8), Inches(1.45), Inches(11.7), Inches(5.85))
    shape.fill.solid(); shape.fill.fore_color.rgb = CODE_BG
    shape.line.color.rgb = APPLE_GREEN; shape.line.width = Pt(3)
    tf = shape.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.25)
    tf.margin_right = Inches(0.25)
    tf.margin_top = Inches(0.18)
    tf.margin_bottom = Inches(0.12)

    def code_line(tf, parts, first_line=False):
        """parts = [(text, color, bold), ...]"""
        p = tf.paragraphs[0] if first_line else tf.add_paragraph()
        for text, color, bold in parts:
            run(p, text, "Consolas", 9.5, bold, color)
        p.space_after = Pt(2)
        p.line_spacing = Pt(13)
        return p

    code_line(tf, [
        ("Goal ", CODE_KW, True),
        ("Collect all required request fields first, then run the AI/WorkIQ overlap", CODE_TEXT, False),
    ], first_line=True)
    code_line(tf, [
        ("check, then proceed based on user choice, and only then trigger the", CODE_TEXT, False),
    ])
    code_line(tf, [
        ("Power Automate creation flow. ", CODE_TEXT, False),
        ("Instructions", CODE_KW, True),
    ])
    code_line(tf, [("", CODE_TEXT, False)])  # blank line

    code_line(tf, [
        ("Intake first (collect structured fields) ", CODE_KW, True),
        ("DO NOT DIRECTLY ASK values", CODE_TEXT, False),
    ])
    code_line(tf, [
        ('to the user, but MAKE SURE TO LAUNCH the "New project request" topic.,', CODE_TEXT, False),
    ])
    code_line(tf, [
        ("it asks the user the required fields for the project request (Title,", CODE_TEXT, False),
    ])
    code_line(tf, [
        ("Requester, Sponsor, Objective/Outcome, Benefits, Urgency, Deadline,", CODE_TEXT, False),
    ])
    code_line(tf, [
        ("Impacts, Dependencies/Constraints, Effort size, Main risks, Notes).", CODE_TEXT, False),
    ])
    code_line(tf, [
        ("Confirm you have captured all fields inside each expected global variable", CODE_TEXT, False),
    ])
    code_line(tf, [
        ("correctly (or ask follow-ups if something is missing).", CODE_TEXT, False),
    ])
    code_line(tf, [("", CODE_TEXT, False)])  # blank line

    code_line(tf, [
        ("Then run AI search ", CODE_KW, True),
        ('(Work IQ copilot_chat MCP tool using this exact query', CODE_TEXT, False),
    ])
    code_line(tf, [
        ('"Search inside my OneDrive folder', CODE_STR, False),
    ])
    code_line(tf, [
        ("'CopilotStudioPOCs/POC_ProjectIntake/KnowledgeBase' for documents", CODE_STR, False),
    ])
    code_line(tf, [
        ('related to improving customer data quality in CRM and reducing', CODE_STR, False),
    ])
    code_line(tf, [
        ('duplicate customer records.")', CODE_STR, False),
        (" to detect overlap and dependencies.", CODE_KW, True),
    ])
    code_line(tf, [("", CODE_TEXT, False)])

    code_line(tf, [
        ("Search ONLY STRICTLY inside my personal OneDrive KnowledgeBase folder", CODE_TEXT, False),
    ])
    code_line(tf, [
        ("for similar initiatives, related project documents, risks, stakeholders,", CODE_TEXT, False),
    ])
    code_line(tf, [
        ("dependencies, or governance notes ", CODE_TEXT, False),
        ("READING THE FILES CONTENT.", CODE_KW, True),
    ])
    code_line(tf, [
        ("Prefer concise results (top 3–5). Include full direct URI to file", CODE_TEXT, False),
    ])
    code_line(tf, [
        ("sources in the response text AND SUMMARIZE THEIR CONTENT AFTER OPENING", CODE_TEXT, False),
    ])
    code_line(tf, [
        ("AND READING IT, mentioning any stakeholder PERSON FULL NAME if present.", CODE_TEXT, False),
    ])
    code_line(tf, [
        ("DO NOT MAKE UP NON-EXISTING DOCUMENT LINKS.", CODE_KW, True),
    ])

    # ── Slide B: Prompt Part 2 (Present findings + Power Automate) ──
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(s, "SYSTEM INSTRUCTIONS (2/2)", "PROMPT VERBATIM — COPILOT STUDIO · GPT-5 CHAT")
    green_line(s)

    shape = s.shapes.add_shape(1, Inches(0.8), Inches(1.45), Inches(11.7), Inches(5.85))
    shape.fill.solid(); shape.fill.fore_color.rgb = CODE_BG
    shape.line.color.rgb = APPLE_GREEN; shape.line.width = Pt(3)
    tf = shape.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.25)
    tf.margin_right = Inches(0.25)
    tf.margin_top = Inches(0.18)
    tf.margin_bottom = Inches(0.12)

    code_line(tf, [
        ("Present findings and ask for next step ", CODE_KW, True),
        ("If relevant context is found,", CODE_TEXT, False),
    ], first_line=True)
    code_line(tf, [
        ("present: Top related initiatives/documents (with link) Possible", CODE_TEXT, False),
    ])
    code_line(tf, [
        ("overlaps Risks/dependencies/stakeholders to consider.", CODE_TEXT, False),
    ])
    code_line(tf, [
        ("DO NOT REPEAT THE RESPONSES GATHERED FROM THE USER HERE.", CODE_KW, True),
    ])
    code_line(tf, [("", CODE_TEXT, False)])

    code_line(tf, [
        ('Then ask the user, with the exact question "§§§ What would you like', CODE_TEXT, False),
    ])
    code_line(tf, [
        ('to do next?", to choose one option:', CODE_TEXT, False),
    ])
    code_line(tf, [
        ("  ✅ Proceed as-is", CODE_STR, False),
    ])
    code_line(tf, [
        ("  ⟲  Revise the scope ", CODE_STR, False),
        ('(ask what <value to change>, then launch', CODE_TEXT, False),
    ])
    code_line(tf, [
        ('     the related "Revise Project Request <value to change>" topic', CODE_TEXT, False),
    ])
    code_line(tf, [
        ('     accordingly; you may re-run the search if needed).', CODE_TEXT, False),
    ])
    code_line(tf, [("", CODE_TEXT, False)])

    code_line(tf, [
        ("Ensure that all other captured fields exactly as they stand", CODE_TEXT, False),
    ])
    code_line(tf, [
        ('(do not omit, and do not replace "none" with null).', CODE_TEXT, False),
    ])
    code_line(tf, [("", CODE_TEXT, False)])

    code_line(tf, [
        ("Only after user confirmation, and AI Notes variable is set", CODE_KW, True),
    ])
    code_line(tf, [
        ("then trigger Power Automate. ", CODE_KW, True),
        ('Trigger the Power Automate "Create new', CODE_TEXT, False),
    ])
    code_line(tf, [
        ('project request" flow to create the tracker entry and generate the', CODE_TEXT, False),
    ])
    code_line(tf, [
        ("official one-pager, using the latest captured field values (including", CODE_TEXT, False),
    ])
    code_line(tf, [
        ("any revised fields and/or AI Notes updating the related global", CODE_TEXT, False),
    ])
    code_line(tf, [
        ("variables), otherwise using exact user provided answers as input", CODE_TEXT, False),
    ])
    code_line(tf, [
        ('parameters i.e. do not convert "none" to null, as all the parameters', CODE_TEXT, False),
    ])
    code_line(tf, [
        ("are required).", CODE_TEXT, False),
    ])


def slide_prompt_screenshot(prs):
    """Instructions screenshot side-by-side with key points."""
    s = prs.slides.add_slide(prs.slide_layouts[6])

    # Left: image (55% width)
    img_w = Inches(7.2); img_h = Inches(6.2)
    img_fit(s, SCREENSHOTS / "prompt_detail_top.png",
            Inches(0.5), Inches(0.6), img_w, img_h, border=True)

    # Right: text panel
    x = Inches(8.2); w = Inches(4.7)

    tf = tb(s, x, Inches(1.2), w, Inches(0.3))
    para(tf, "COPILOT STUDIO — OVERVIEW", "Arial Black", 10, True,
         BRIGHT_GREEN, PP_ALIGN.LEFT, Pt(0), first=True)

    tf = tb(s, x, Inches(1.7), w, Inches(0.7))
    para(tf, "INSTRUCTIONS\nDELL'AGENTE", "Arial Black", 20, True,
         BLACK, PP_ALIGN.LEFT, Pt(8), first=True)

    tf = tb(s, x, Inches(2.7), w, Inches(0.5))
    para(tf, "Workflow in 4 fasi sequenziali:", "Arial", 13, False,
         BLACK, PP_ALIGN.LEFT, Pt(8), first=True)

    tf = tb(s, x, Inches(3.4), w, Inches(3.5))
    bullet(tf, "Intake: raccolta di 12 campi strutturati via topic dedicato", 12)
    tf.paragraphs[0].text = ""  # remove default blank
    bullet(tf, "AI Search: query esatta su KnowledgeBase OneDrive via Work IQ MCP", 12)
    bullet(tf, "Decisione: presentazione risultati + scelta utente (proceed / revise)", 12)
    bullet(tf, "Creazione: Power Automate genera tracker entry + OnePager SharePoint", 12)

    tf = tb(s, x, Inches(5.7), w, Inches(0.4))
    para(tf, "Modello: GPT-5 Chat", "Arial", 11, False,
         MID_GRAY, PP_ALIGN.LEFT, Pt(0), first=True)

    # Green accent stripe on left edge of text panel
    shape = s.shapes.add_shape(1, x - Inches(0.15), Inches(1.0), Pt(4), Inches(5.2))
    shape.fill.solid(); shape.fill.fore_color.rgb = APPLE_GREEN
    shape.line.fill.background()


def slide_user_prompt(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(s, "PROMPT UTENTE", "IL PROMPT CHE AVVIA IL PROCESSO")
    green_line(s)

    # Quote box
    shape = s.shapes.add_shape(1, Inches(2.2), Inches(1.7), Inches(8.9), Inches(1.4))
    shape.fill.solid(); shape.fill.fore_color.rgb = CODE_BG
    shape.line.color.rgb = VIVID_BLUE; shape.line.width = Pt(3)
    tf = shape.text_frame; tf.word_wrap = True
    tf.margin_left = Inches(0.25)
    tf.margin_right = Inches(0.25)
    tf.margin_top = Inches(0.18)
    p = tf.paragraphs[0]
    run(p, '"I want to open a new internal project request to improve\n'
           'customer data quality in CRM and reduce duplicate\n'
           'customer records"', "Consolas", 15, False, CODE_STR)
    p.alignment = PP_ALIGN.CENTER
    p.line_spacing = Pt(22)

    # Screenshot (centered, large)
    img_fit(s, SCREENSHOTS / "03_prompt_user_request.png",
            Inches(1.5), Inches(3.5), Inches(10.3), Inches(3.2), border=True)

    tf = tb(s, Inches(2.2), Inches(6.85), Inches(8.9), Inches(0.4))
    para(tf, 'L\'agente interpreta il prompt e instrada automaticamente al topic "New project request"',
         "Arial", 10, True, DARK_GRAY, PP_ALIGN.CENTER, Pt(0), first=True)


def slide_image_side(prs, img, evidence, title, body, bullets_list, footer=None):
    """Screenshot left (55%) + text right (45%)."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    img_w = Inches(6.8)

    # Image
    img_fit(s, SCREENSHOTS / img, Inches(0.5), Inches(0.45), img_w, Inches(6.5), border=True)

    # Vertical green accent
    x = Inches(7.55)
    shape = s.shapes.add_shape(1, x, Inches(1.0), Pt(4), Inches(5.3))
    shape.fill.solid(); shape.fill.fore_color.rgb = APPLE_GREEN
    shape.line.fill.background()

    # Text panel
    tx = Inches(7.95); tw = Inches(5.0)

    y = Inches(1.2)
    tf = tb(s, tx, y, tw, Inches(0.3))
    para(tf, evidence, "Arial Black", 10, True, BRIGHT_GREEN, PP_ALIGN.LEFT, Pt(0), first=True)
    y += Inches(0.5)

    tf = tb(s, tx, y, tw, Inches(0.8))
    para(tf, title, "Arial Black", 18, True, BLACK, PP_ALIGN.LEFT, Pt(8), first=True)
    y += Inches(0.9)

    if body:
        tf = tb(s, tx, y, tw, Inches(0.7))
        para(tf, body, "Arial", 12, False, BLACK, PP_ALIGN.LEFT, Pt(10), first=True)
        y += Inches(0.7)

    if bullets_list:
        tf = tb(s, tx, y, tw, Inches(3.2))
        tf.paragraphs[0].text = ""
        for b in bullets_list:
            bullet(tf, b, 11)
        y += Inches(0.4 + 0.28 * len(bullets_list))

    if footer:
        # Footer with light bg
        shape = s.shapes.add_shape(1, tx, y + Inches(0.25), tw, Inches(0.7))
        shape.fill.solid(); shape.fill.fore_color.rgb = LIGHT_GRAY
        shape.line.fill.background()
        ftf = shape.text_frame; ftf.word_wrap = True
        ftf.margin_left = Inches(0.12)
        ftf.margin_top = Inches(0.1)
        para(ftf, footer, "Arial", 10, False, DARK_GRAY, PP_ALIGN.LEFT, Pt(0), first=True)


def slide_tool_table(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(s, "MAPPA DEI TOOL", "TUTTI I COMPONENTI COINVOLTI NEL FLUSSO")
    green_line(s)

    rows = 6; cols = 4
    tbl = s.shapes.add_table(rows, cols, Inches(1.0), Inches(1.8),
                             Inches(11.3), Inches(4.6)).table
    tbl.columns[0].width = Inches(3.3)
    tbl.columns[1].width = Inches(2.0)
    tbl.columns[2].width = Inches(4.5)
    tbl.columns[3].width = Inches(1.5)

    for i, h in enumerate(["TOOL", "TIPO", "FUNZIONE", "DURATA"]):
        c = tbl.cell(0, i); c.text = h
        c.fill.solid(); c.fill.fore_color.rgb = APPLE_GREEN
        c.margin_left = Inches(0.08)
        c.margin_top = Inches(0.04)
        c.margin_bottom = Inches(0.04)
        p = c.text_frame.paragraphs[0]
        p.font.name = "Arial Black"; p.font.size = Pt(11)
        p.font.color.rgb = WHITE; p.font.bold = True

    data = [
        ("copilot_chat", "MCP", "AI Search nella KnowledgeBase OneDrive", "55.31s"),
        ("New project request", "Topic", "Raccolta 12 campi strutturati (intake form)", "~3 min"),
        ("Revise Project Request Title", "Topic", "Revisione titolo su richiesta utente", "~15s"),
        ("Create new project request", "Power Automate", "Genera OnePager su SharePoint + Request ID", "10.11s"),
        ("Copilot in Word", "M365 Copilot", "Formattazione AI Notes da markdown a Word", "~20s"),
    ]
    type_colors = {"MCP": VIVID_BLUE, "Topic": APPLE_GREEN,
                   "Power Automate": AMBER, "M365 Copilot": AMBER}

    for r, (tool, tipo, func, dur) in enumerate(data, 1):
        for c_idx, val in enumerate([tool, tipo, func, dur]):
            cell = tbl.cell(r, c_idx); cell.text = val
            cell.margin_left = Inches(0.08)
            cell.margin_right = Inches(0.05)
            cell.margin_top = Inches(0.04)
            cell.margin_bottom = Inches(0.04)
            p = cell.text_frame.paragraphs[0]
            p.font.name = "Arial"; p.font.size = Pt(12); p.font.color.rgb = BLACK
            if c_idx == 0: p.font.bold = True
            if c_idx == 1:
                p.font.color.rgb = type_colors.get(tipo, BLACK); p.font.bold = True
            if r % 2 == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = LIGHT_GRAY


def slide_dual_screenshot(prs, title, subtitle, img1, cap1, img2, cap2):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(s, title, subtitle)
    green_line(s)

    iw = Inches(5.4); ih = Inches(4.2)
    img_fit(s, SCREENSHOTS / img1, Inches(0.7), Inches(1.8), iw, ih, border=True)
    img_fit(s, SCREENSHOTS / img2, Inches(7.2), Inches(1.8), iw, ih, border=True)

    tf = tb(s, Inches(0.7), Inches(6.2), iw, Inches(0.4))
    para(tf, cap1, "Arial", 10, True, DARK_GRAY, PP_ALIGN.CENTER, Pt(0), first=True)
    tf = tb(s, Inches(7.2), Inches(6.2), iw, Inches(0.4))
    para(tf, cap2, "Arial", 10, True, DARK_GRAY, PP_ALIGN.CENTER, Pt(0), first=True)

    # Vertical separator
    shape = s.shapes.add_shape(1, Inches(6.65), Inches(1.8), Pt(2), ih)
    shape.fill.solid(); shape.fill.fore_color.rgb = RGBColor(0xE0,0xE0,0xE0)
    shape.line.fill.background()


def slide_video(prs, title, subtitle, poster, clip_file, clip_label, duration, extra_content=None):
    """Video phase slide with embedded MP4 video."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    title_bar(s, title, subtitle)
    green_line(s)

    # Green video label bar
    shape = s.shapes.add_shape(1, Inches(1.0), Inches(1.55), Inches(11.3), Inches(0.5))
    shape.fill.solid(); shape.fill.fore_color.rgb = APPLE_GREEN
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = Inches(0.25)
    p = tf.paragraphs[0]
    run(p, f"▶  {clip_label}", "Arial Black", 12, True, WHITE)
    run(p, f"    {duration}", "Arial", 12, False, RGBColor(0xDD,0xFF,0xDD))
    p.alignment = PP_ALIGN.LEFT

    # Embedded video with poster frame
    vid_y = Inches(2.3)
    vid_h = Inches(4.3) if not extra_content else Inches(3.3)
    video_path = CLIPS / clip_file
    embed_video(s, prs, video_path, SCREENSHOTS / poster,
                Inches(1.0), vid_y, Inches(11.3), vid_h)

    if extra_content:
        extra_content(s)

    # Clip reference footnote
    tf = tb(s, Inches(2.5), Inches(7.0), Inches(8.3), Inches(0.35))
    para(tf, f"📎 {clip_file}  ·  {duration}",
         "Arial", 9, False, MID_GRAY, PP_ALIGN.CENTER, Pt(0), first=True)


def slide_statement(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])

    tf = tb(s, Inches(1.8), Inches(1.8), Inches(9.7), Inches(3.2))
    p = tf.paragraphs[0]
    run(p, "DALL'IDEA AL\n", "Arial Black", 44, True, BLACK)
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(12)

    p2 = tf.add_paragraph()
    run(p2, "DOCUMENTO UFFICIALE", "Arial Black", 44, True, APPLE_GREEN)
    p2.alignment = PP_ALIGN.CENTER
    p2.space_after = Pt(12)

    p3 = tf.add_paragraph()
    run(p3, "\nIN 13 MINUTI", "Arial Black", 44, True, BLACK)
    p3.alignment = PP_ALIGN.CENTER

    tf2 = tb(s, Inches(2.8), Inches(5.5), Inches(7.7), Inches(1.2))
    para(tf2, "Un singolo prompt attiva intake strutturato, AI search nella knowledge base "
              "aziendale, e generazione automatica su SharePoint — senza uscire dalla chat.",
         "Arial", 15, False, DARK_GRAY, PP_ALIGN.CENTER, Pt(0), first=True)


def slide_thankyou(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg_gradient(s, DARK_GREEN, APPLE_GREEN)

    tf = tb(s, Inches(1), Inches(2.7), Inches(11.3), Inches(1.5))
    para(tf, "THANK YOU", "Arial Black", 60, True, WHITE, PP_ALIGN.CENTER, Pt(16), first=True)

    tf = tb(s, Inches(2.5), Inches(4.8), Inches(8.3), Inches(0.6))
    para(tf, "Project Request Tracker — Copilot Studio + Work IQ MCP",
         "Arial", 16, False, RGBColor(0xDD,0xDD,0xDD), PP_ALIGN.CENTER, Pt(0), first=True)


# ── Main ─────────────────────────────────────────────────────────────

def main():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # 1. Cover
    slide_cover(prs)

    # 2. Agenda
    slide_agenda(prs)

    # ── SEZIONE: OVERVIEW ──
    # 3. Section: Overview
    slide_section(prs, "🎯", "OVERVIEW",
        "Un agente Copilot Studio che gestisce l'intero ciclo di creazione project request — "
        "dall'intake conversazionale alla generazione automatica di un OnePager su SharePoint.")

    # 4. Flow + KPI
    slide_flow_kpi(prs)

    # ── SEZIONE: PROMPT ──
    # 5. Section: Prompt
    slide_section(prs, "📝", "SYSTEM PROMPT",
        "Le istruzioni di sistema definiscono il comportamento dell'agente: "
        "raccolta dati, AI search, decisione utente e trigger Power Automate.")

    # 6-7. System Instructions verbatim (2 slides)
    slide_prompt_verbatim(prs)

    # 8. Instructions screenshot
    slide_prompt_screenshot(prs)

    # 9. User prompt
    slide_user_prompt(prs)

    # ── SEZIONE: TOOL ──
    # 10. Section: Tool
    slide_section(prs, "🔧", "TOOL UTILIZZATI",
        "Work IQ MCP, Topics Copilot Studio, AI Tool Node e Power Automate "
        "orchestrano il flusso end-to-end.")

    # 11. Work IQ MCP panel
    slide_image_side(prs,
        "04_tool_workiq_mcp_panel.png",
        "MODEL CONTEXT PROTOCOL",
        "WORK IQ COPILOT\n(PREVIEW)",
        "Server MCP che espone il tool copilot_chat per interrogare Microsoft 365:",
        [
            "Ricerca semantica su documenti, email, chat, siti",
            "Accesso a OneDrive, SharePoint, Teams, Mail",
            "Priorità ai tool workload-specific quando il contesto è chiaro",
        ],
        'Query: "Search inside my OneDrive folder \'KnowledgeBase\' for documents related to CRM data quality..."')

    # 12. Topics Editor (moved here — tool authoring context)
    slide_image_side(prs,
        "07_tool_topics_editor.png",
        "COPILOT STUDIO — AUTHORING",
        "TOPICS EDITOR",
        'Il topic "New project request" definisce il flusso di raccolta dati con nodi Question:',
        [
            "Urgency → Global.varUrgency",
            "Expected Effort → Global.varEffort",
            "Deadline → Global.varDeadline",
        ],
        "Ogni variabile globale viene passata al flusso Power Automate per la generazione del OnePager.")

    # 13. Tool table
    slide_tool_table(prs)

    # 14. copilot_chat before/after
    slide_dual_screenshot(prs,
        "COPILOT_CHAT — ESECUZIONE",
        "55 SECONDI DI AI SEARCH NELLA KNOWLEDGEBASE",
        "05_tool_copilot_chat_working.png", 'Tool in esecuzione — stato "Working"',
        "06_tool_copilot_chat_completed.png", "Tool completato — risultati disponibili")

    # ── SEZIONE: DEMO VIDEO ──
    # 15. Section: Demo
    slide_section(prs, "🎬", "DEMO VIDEO",
        "Registrazioni del flusso operativo completo: dal prompt iniziale "
        "alla generazione del documento finale.")

    # 16. Video: Setup
    slide_video(prs, "FASE 1 — SETUP AGENTE",
        "CONFIGURAZIONE COPILOT STUDIO: OVERVIEW, INSTRUCTIONS, TOOLS",
        "01_prompt_system_instructions.png",
        "clip1_setup_agent_instructions.mp4", "Setup Agente & Instructions", "3:40")

    # 17. Video: Prompt + Data
    slide_video(prs, "FASE 2 — PROMPT & RACCOLTA DATI",
        "L'UTENTE INVIA IL PROMPT E L'AGENTE RACCOGLIE I 12 CAMPI",
        "03_prompt_user_request.png",
        "clip2_user_prompt_data_collection.mp4", "Prompt Utente & Data Collection", "2:40")

    # 18. Video: AI Search
    slide_video(prs, "FASE 3 — AI SEARCH CON WORK IQ",
        "IL TOOL COPILOT_CHAT CERCA NELLA KNOWLEDGEBASE",
        "10_demo_ai_search_results.png",
        "clip3_ai_search_workiq_results.mp4", "AI Search & Risultati KnowledgeBase", "2:00")

    # 19. AI Results detail
    slide_image_side(prs,
        "10_demo_ai_search_results.png",
        "TOP RELATED DOCUMENTS",
        "RISULTATI AI SEARCH",
        None,
        [
            "Customer_Master_Cleanup_Initiative.docx — Strategic initiative, duplicati, naming, VAT IDs",
            "CRM_Data_Quality_Meeting_Notes.docx — Stakeholder discussions, validation rules",
            "Customer_Data_Risk_Assessment.docx — Risk analysis, CRM/ERP sync failures",
        ],
        "⚠️ Temi emergenti: governance già discussa, rischi ERP sync, stakeholder Sales Ops/IT/Data Governance")

    # 20. Knowledge Base (moved here — source docs context after AI results)
    slide_image_side(prs,
        "15_demo_knowledge_source_doc.png",
        "KNOWLEDGE BASE",
        "DOCUMENTI SORGENTE",
        "I documenti OneDrive costituiscono il corpus per la ricerca semantica del tool copilot_chat:",
        [
            "Customer Master Cleanup Initiative — Piano strategico pulizia dati CRM",
            "CRM Data Quality Meeting Notes — Verbali riunioni stakeholder",
            "Customer Data Risk Assessment — Analisi rischi e mitigazioni",
        ],
        "L'agente cerca solo nella cartella specificata, garantendo risultati pertinenti e controllati.")

    # 21. Decision Point (NEW — screenshot 11)
    slide_image_side(prs,
        "11_demo_decision_point.png",
        "DECISION POINT",
        "§§§ WHAT WOULD YOU\nLIKE TO DO NEXT?",
        "L'agente presenta i risultati della ricerca e chiede all'utente di scegliere:",
        [
            "✅ Proceed as-is — conferma e procedi alla creazione",
            "⟲ Revise the scope — modifica un campo e riesegui la ricerca se necessario",
        ],
        "L'utente sceglie di rivedere il titolo prima di procedere con la creazione.")

    # 22. Revise Title Flow (NEW — screenshot 08)
    slide_image_side(prs,
        "08_tool_revise_title_flow.png",
        "COPILOT STUDIO — TOPIC FLOW",
        "REVISE PROJECT\nREQUEST TITLE",
        "Il topic dedicato gestisce la revisione del titolo della request:",
        [
            "Nodo Question per il nuovo valore del titolo",
            "Aggiornamento di Global.varTitle con il valore rivisto",
            "Ritorno al flusso principale per la conferma finale",
        ],
        None)

    # 23. Video: Revise & Create
    def revise_extra(s):
        shape = s.shapes.add_shape(1, Inches(1.0), Inches(6.2), Inches(5.4), Inches(0.7))
        shape.fill.solid(); shape.fill.fore_color.rgb = LIGHT_GRAY
        shape.line.fill.background()
        ftf = shape.text_frame; ftf.word_wrap = True
        ftf.margin_left = Inches(0.15)
        ftf.margin_top = Inches(0.08)
        para(ftf, "PRIMA", "Arial Black", 9, True, BRIGHT_GREEN, PP_ALIGN.LEFT, Pt(4), first=True)
        para(ftf, 'Title: "Improve CRM data quality..."', "Arial", 11, False, DARK_GRAY)

        shape2 = s.shapes.add_shape(1, Inches(6.9), Inches(6.2), Inches(5.4), Inches(0.7))
        shape2.fill.solid(); shape2.fill.fore_color.rgb = LIGHT_GRAY
        shape2.line.fill.background()
        ftf2 = shape2.text_frame; ftf2.word_wrap = True
        ftf2.margin_left = Inches(0.15)
        ftf2.margin_top = Inches(0.08)
        para(ftf2, "DOPO REVISIONE", "Arial Black", 9, True, BRIGHT_GREEN, PP_ALIGN.LEFT, Pt(4), first=True)
        para(ftf2, 'Title: "CRM data deduplication"', "Arial", 11, True, BLACK)

    slide_video(prs, "FASE 4 — REVISIONE & CREAZIONE",
        "L'UTENTE RIVEDE IL TITOLO, CONFERMA, E POWER AUTOMATE GENERA LA REQUEST",
        "09_tool_create_request_completed.png",
        "clip4_revise_create_request.mp4", "Revisione Titolo & Creazione", "1:10",
        extra_content=revise_extra)

    # 24. Video: OnePager
    slide_video(prs, "FASE 5 — ONEPAGER SU SHAREPOINT",
        "DOCUMENTO GENERATO + FORMATTAZIONE CON COPILOT IN WORD",
        "14_demo_final_onepager.png",
        "clip5_onepager_copilot_word.mp4", "OnePager & Copilot Word Formatting", "2:43")

    # 25. Copilot Word Formatting (NEW — screenshot 13)
    slide_image_side(prs,
        "13_demo_copilot_word_formatting.png",
        "M365 COPILOT IN WORD",
        "FORMATTAZIONE\nAI NOTES",
        "Copilot in Word converte le AI Notes da markdown grezzo a formato professionale:",
        [
            "Prompt: 'Format this content with proper headings and structure'",
            "Intestazioni, tabelle e sezioni generate automaticamente",
            "Documento finale pronto per la condivisione su SharePoint",
        ],
        None)

    # 26. Before/After
    slide_dual_screenshot(prs,
        "ONEPAGER — PRIMA E DOPO",
        "DA MARKDOWN RAW A DOCUMENTO FORMATTATO CON COPILOT IN WORD",
        "12_demo_onepager_raw.png", "AI Notes in formato markdown grezzo",
        "14_demo_final_onepager.png", "Documento finale formattato con Copilot Word")

    # ── CHIUSURA ──
    # 27. Statement
    slide_statement(prs)

    # 28. Thank You
    slide_thankyou(prs)

    prs.save(str(OUTPUT))
    print(f"✅ PPTX generato: {OUTPUT}")
    print(f"   Slide: {len(prs.slides)}")
    sz = OUTPUT.stat().st_size / 1024 / 1024
    print(f"   Size: {sz:.1f} MB")


if __name__ == "__main__":
    main()
