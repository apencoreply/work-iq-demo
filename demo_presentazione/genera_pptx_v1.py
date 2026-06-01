#!/usr/bin/env python3
"""
Genera due file PPTX dalla presentazione Work IQ Demo Comparativa:
  1. presentazione.pptx          — versione testuale (senza immagini)
  2. presentazione_visual.pptx   — versione con screenshot integrati

Richiede: pip install python-pptx Pillow
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Paths ──
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
SCREENSHOTS_DIR = os.path.join(SCRIPT_DIR, "screenshots")

# ── Brand Colors ──
GREEN       = RGBColor(0x76, 0xB8, 0x2A)
DARK_GREEN  = RGBColor(0x2D, 0x6A, 0x2E)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x1A, 0x1A, 0x1A)
GRAY        = RGBColor(0x55, 0x55, 0x55)
LIGHT_GRAY  = RGBColor(0xF5, 0xF5, 0xF5)
BG_DARK     = RGBColor(0x0D, 0x11, 0x17)
RED         = RGBColor(0xEF, 0x44, 0x44)
YELLOW      = RGBColor(0xF5, 0x9E, 0x0B)
PURPLE      = RGBColor(0x8B, 0x5C, 0xF6)
BRIGHT_GREEN = RGBColor(0x8C, 0xC6, 0x3F)

SLIDE_WIDTH  = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)


# ═══════════════════════════════════════════════════════
# Utility helpers
# ═══════════════════════════════════════════════════════

def set_slide_bg(slide, color):
    """Set solid background color for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def set_slide_gradient(slide, color1, color2):
    """Set a two-stop gradient background."""
    bg = slide.background
    fill = bg.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = color1
    fill.gradient_stops[0].position = 0.0
    fill.gradient_stops[1].color.rgb = color2
    fill.gradient_stops[1].position = 1.0


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=BLACK, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="Arial"):
    """Add a textbox with a single text run."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font_name
    return txBox, tf


def add_rich_textbox(slide, left, top, width, height):
    """Add a textbox and return the text_frame for multi-paragraph use."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    return txBox, tf


def add_run(paragraph, text, font_size=14, color=BLACK, bold=False, italic=False, font_name="Arial"):
    """Add a formatted run to a paragraph."""
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font_name
    return run


def add_paragraph(tf, text="", font_size=14, color=BLACK, bold=False,
                  italic=False, alignment=PP_ALIGN.LEFT, font_name="Arial",
                  space_after=Pt(6)):
    """Add a new paragraph with a single run to a text_frame."""
    p = tf.add_paragraph()
    p.alignment = alignment
    p.space_after = space_after
    if text:
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.italic = italic
        run.font.name = font_name
    return p


def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None):
    """Add a rounded rectangle shape."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(2)
    else:
        shape.line.fill.background()
    return shape


def add_image_centered(slide, img_path, max_width, max_height, top_offset):
    """Add an image centered horizontally, fitting within max bounds."""
    from PIL import Image as PILImage
    with PILImage.open(img_path) as img:
        img_w, img_h = img.size
    aspect = img_w / img_h
    target_w = max_width
    target_h = int(target_w / aspect)
    if target_h > max_height:
        target_h = max_height
        target_w = int(target_h * aspect)
    left = (SLIDE_WIDTH - target_w) // 2
    return slide.shapes.add_picture(img_path, left, top_offset, target_w, target_h)


# ═══════════════════════════════════════════════════════
#  TEXT-ONLY PRESENTATION  (presentazione.pptx)
# ═══════════════════════════════════════════════════════

def build_text_pptx():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank = prs.slide_layouts[6]  # blank layout

    # ── SLIDE 1: COVER ──
    sl = prs.slides.add_slide(blank)
    set_slide_gradient(sl, DARK_GREEN, GREEN)
    add_textbox(sl, Inches(1), Inches(3.5), Inches(9), Inches(2),
                "WORK IQ\nDEMO\nCOMPARATIVA", 44, WHITE, True, font_name="Arial Black")
    add_textbox(sl, Inches(1), Inches(5.5), Inches(9), Inches(0.5),
                "4 livelli di intelligenza — Stesso prompt, risposte diverse", 16, WHITE)
    add_textbox(sl, Inches(1), Inches(6.1), Inches(9), Inches(0.4),
                "1 Giugno 2026 · Copilot Studio + Work IQ MCP", 13, RGBColor(0xCC, 0xCC, 0xCC))

    # ── SLIDE 2: AGENDA ──
    sl = prs.slides.add_slide(blank)
    add_textbox(sl, Inches(0), Inches(0.6), SLIDE_WIDTH, Inches(0.6),
                "AGENDA", 28, BLACK, True, PP_ALIGN.CENTER, "Arial Black")
    agenda_items = [
        "1  La sfida: perché un LLM generico non basta",
        "2  Setup: 4 agenti a confronto",
        "3  Prompt 1 — Preparazione alla riunione",
        "4  Prompt 2 — Riassunto email non lette",
        "5  Prompt 3 — Azione: invio email",
        "6  Riepilogo e messaggi chiave",
    ]
    _, tf = add_rich_textbox(sl, Inches(3), Inches(1.8), Inches(7), Inches(5))
    tf.paragraphs[0].space_after = Pt(0)
    for item in agenda_items:
        p = add_paragraph(tf, font_size=18, space_after=Pt(14))
        num, rest = item.split("  ", 1)
        add_run(p, num + "  ", 22, GREEN, True, font_name="Arial Black")
        add_run(p, rest, 18, BLACK)

    # ── SLIDE 3: LA SFIDA ──
    sl = prs.slides.add_slide(blank)
    add_textbox(sl, Inches(0), Inches(2.5), SLIDE_WIDTH, Inches(1),
                "LA SFIDA", 36, GREEN, True, PP_ALIGN.CENTER, "Arial Black")
    add_textbox(sl, Inches(2), Inches(3.8), Inches(9.3), Inches(1.5),
                "Un modello linguistico senza contesto aziendale è come un consulente brillante che non ha mai visto i tuoi dati.",
                20, GRAY, False, PP_ALIGN.CENTER)

    # ── SLIDE 4: STATEMENT ──
    sl = prs.slides.add_slide(blank)
    set_slide_bg(sl, BLACK)
    add_textbox(sl, Inches(1), Inches(2), Inches(11), Inches(3.5),
                "SENZA CONTESTO,\nL'AI NON SERVE\nIN AZIENDA.", 40, WHITE, True, PP_ALIGN.CENTER, "Arial Black")
    # Highlight "NON SERVE" in green via a separate box on top
    add_textbox(sl, Inches(4.5), Inches(3.35), Inches(4.5), Inches(0.8),
                "NON SERVE", 40, GREEN, True, PP_ALIGN.CENTER, "Arial Black")

    # ── SLIDE 5: 4 LIVELLI ──
    sl = prs.slides.add_slide(blank)
    add_textbox(sl, Inches(0), Inches(0.4), SLIDE_WIDTH, Inches(0.6),
                "I 4 LIVELLI DI INTELLIGENZA", 26, BLACK, True, PP_ALIGN.CENTER, "Arial Black")
    add_textbox(sl, Inches(0), Inches(1.0), SLIDE_WIDTH, Inches(0.4),
                "STESSO PROMPT, CAPACITÀ CRESCENTI", 13, GRAY, True, PP_ALIGN.CENTER, "Arial Black")

    cards = [
        ("1", "BASE", "Solo LLM generico.\nNessun tool, nessun dato.", "—", RED),
        ("2", "CONNESSIONI", "Accesso diretto a dati\nMicrosoft 365.", "Mail · Calendar · User", YELLOW),
        ("3", "WORK IQ COPILOT", "Semantic Index: capisce\nil contesto nel tempo.", "Semantic Index · Memory\nCross-app", PURPLE),
        ("4", "COMPLETO", "Intelligenza profonda +\naccesso diretto.\nCapisce e agisce.", "Copilot + Mail +\nCalendar + User", GREEN),
    ]
    card_w = Inches(2.8)
    gap = Inches(0.3)
    total = card_w * 4 + gap * 3
    start_x = (SLIDE_WIDTH - total) // 2
    for i, (num, name, desc, tools, color) in enumerate(cards):
        x = start_x + i * (card_w + gap)
        y = Inches(1.8)
        h = Inches(4.5)
        rect = add_rounded_rect(sl, x, y, card_w, h, LIGHT_GRAY, color)
        # Top color bar
        bar = slide_shapes_add_shape = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, card_w, Inches(0.15))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        # Number
        add_textbox(sl, x, y + Inches(0.4), card_w, Inches(0.7),
                    num, 32, GRAY, True, PP_ALIGN.CENTER, "Arial Black")
        # Name
        add_textbox(sl, x, y + Inches(1.2), card_w, Inches(0.5),
                    name, 14, BLACK, True, PP_ALIGN.CENTER, "Arial Black")
        # Desc
        add_textbox(sl, x + Inches(0.2), y + Inches(1.9), card_w - Inches(0.4), Inches(1.5),
                    desc, 12, GRAY, False, PP_ALIGN.CENTER)
        # Tools
        add_textbox(sl, x + Inches(0.2), y + Inches(3.5), card_w - Inches(0.4), Inches(0.8),
                    tools, 10, GRAY, False, PP_ALIGN.CENTER)

    # ── SLIDE 6: SECTION PROMPT 1 ──
    sl = prs.slides.add_slide(blank)
    add_textbox(sl, Inches(0), Inches(2.5), SLIDE_WIDTH, Inches(0.8),
                "PROMPT 1", 36, GREEN, True, PP_ALIGN.CENTER, "Arial Black")
    add_textbox(sl, Inches(2), Inches(3.8), Inches(9.3), Inches(1),
                '"Cosa devo sapere prima della mia prossima riunione?"',
                20, GRAY, False, PP_ALIGN.CENTER)

    # ── SLIDE 7: PROMPT 1 — Agent 1 vs 2 ──
    sl = prs.slides.add_slide(blank)
    add_textbox(sl, Inches(0), Inches(0.3), SLIDE_WIDTH, Inches(0.5),
                "PROMPT 1 — RISPOSTE A CONFRONTO", 22, BLACK, True, PP_ALIGN.CENTER, "Arial Black")
    add_textbox(sl, Inches(0), Inches(0.8), SLIDE_WIDTH, Inches(0.4),
                "BASE vs CONNESSIONI", 12, GRAY, True, PP_ALIGN.CENTER, "Arial Black")

    # Left column - Agent 1
    col_w = Inches(5.8)
    left_x = Inches(0.5)
    right_x = Inches(7)

    add_textbox(sl, left_x, Inches(1.5), col_w, Inches(0.4),
                "🔴  1  AGENTE BASE", 14, RED, True)
    rect = add_rounded_rect(sl, left_x, Inches(2.0), col_w, Inches(3.8), LIGHT_GRAY)
    rect.line.color.rgb = RED
    add_textbox(sl, left_x + Inches(0.2), Inches(2.1), col_w - Inches(0.4), Inches(3.5),
                "Non risultano informazioni specifiche sulla tua prossima riunione nei dati disponibili. "
                "Per offrirti un briefing completo, avrei bisogno di dettagli come l'orario, i partecipanti o l'argomento della riunione.",
                14, BLACK)
    add_textbox(sl, left_x, Inches(5.9), col_w, Inches(0.5),
                "🔴 Educato ma inutile. Non sa nulla.", 13, RED, True)

    # Right column - Agent 2
    add_textbox(sl, right_x, Inches(1.5), col_w, Inches(0.4),
                "🟡  2  CONNESSIONI", 14, YELLOW, True)
    rect = add_rounded_rect(sl, right_x, Inches(2.0), col_w, Inches(3.8), LIGHT_GRAY)
    rect.line.color.rgb = YELLOW
    add_textbox(sl, right_x + Inches(0.2), Inches(2.1), col_w - Inches(0.4), Inches(3.5),
                "Horizon 3.0 - Stand up tecnico\n"
                "Oggi, 09:30 – 10:00 · Microsoft Teams\n"
                "Organizzatore: Marco Bianchi\n"
                "Partecipanti: Elena Rossi ✓, Luca Moretti, Sara Conti, Davide Rizzo (in attesa)\n"
                "Tu: accettato provvisoriamente",
                14, BLACK)
    add_textbox(sl, right_x, Inches(5.9), col_w, Inches(0.5),
                "🟡 Vede il calendario. Dati puntuali.", 13, YELLOW, True)

    # ── SLIDE 8: PROMPT 1 — Agent 3 vs 4 ──
    sl = prs.slides.add_slide(blank)
    add_textbox(sl, Inches(0), Inches(0.3), SLIDE_WIDTH, Inches(0.5),
                "PROMPT 1 — RISPOSTE A CONFRONTO", 22, BLACK, True, PP_ALIGN.CENTER, "Arial Black")
    add_textbox(sl, Inches(0), Inches(0.8), SLIDE_WIDTH, Inches(0.4),
                "WORK IQ COPILOT vs COMPLETO", 12, GRAY, True, PP_ALIGN.CENTER, "Arial Black")

    # Agent 3
    add_textbox(sl, left_x, Inches(1.5), col_w, Inches(0.4),
                "🟣  3  WORK IQ COPILOT", 14, PURPLE, True)
    rect = add_rounded_rect(sl, left_x, Inches(2.0), col_w, Inches(3.8), LIGHT_GRAY)
    rect.line.color.rgb = PURPLE
    add_textbox(sl, left_x + Inches(0.2), Inches(2.1), col_w - Inches(0.4), Inches(3.5),
                "📅 Meridian - Supply Chain Status - Weekly\n"
                "Domani 08:30 – 09:00 · Microsoft Teams\n"
                "Organizzatore: Paolo Gentile\n\n"
                "🧩 Meeting di stato settimanale su Meridian / Supply Chain Platform\n"
                "🔍 Parte di un filone di meeting (Phoenix Stand-Up, Horizon SAL, Atlas Weekly) — preparati su rischi e dipendenze",
                13, BLACK)
    add_textbox(sl, left_x, Inches(5.9), col_w, Inches(0.5),
                "🟣 Analisi strutturata, contesto storico, suggerimenti proattivi", 12, PURPLE, True)

    # Agent 4
    add_textbox(sl, right_x, Inches(1.5), col_w, Inches(0.4),
                "🟢  4  COMPLETO", 14, GREEN, True)
    rect = add_rounded_rect(sl, right_x, Inches(2.0), col_w, Inches(3.8), LIGHT_GRAY)
    rect.line.color.rgb = GREEN
    add_textbox(sl, right_x + Inches(0.2), Inches(2.1), col_w - Inches(0.4), Inches(3.5),
                "Horizon 3.0 - Stand up tecnico\n"
                "09:30 – 10:00 · Microsoft Teams\n"
                "Organizzatore: Marco Bianchi\n"
                "Elena Rossi ✓ · Luca, Sara, Davide (in attesa)\n"
                "Tu: tentativo\n\n"
                "💡 Suggerimento: prepara aggiornamenti rapidi su attività Horizon 3.0 ed eventuali blocchi.",
                13, BLACK)
    add_textbox(sl, right_x, Inches(5.9), col_w, Inches(0.5),
                "🟢 Dettagli operativi + stato RSVP + azioni concrete", 12, GREEN, True)

    # ── SLIDE 9: SECTION PROMPT 2 ──
    sl = prs.slides.add_slide(blank)
    add_textbox(sl, Inches(0), Inches(2.5), SLIDE_WIDTH, Inches(0.8),
                "PROMPT 2", 36, GREEN, True, PP_ALIGN.CENTER, "Arial Black")
    add_textbox(sl, Inches(2), Inches(3.8), Inches(9.3), Inches(1),
                '"Riassumi le email non lette più importanti di oggi."',
                20, GRAY, False, PP_ALIGN.CENTER)

    # ── SLIDE 10: PROMPT 2 — Agent 1 vs 2 ──
    sl = prs.slides.add_slide(blank)
    add_textbox(sl, Inches(0), Inches(0.3), SLIDE_WIDTH, Inches(0.5),
                "PROMPT 2 — RISPOSTE A CONFRONTO", 22, BLACK, True, PP_ALIGN.CENTER, "Arial Black")
    add_textbox(sl, Inches(0), Inches(0.8), SLIDE_WIDTH, Inches(0.4),
                "BASE vs CONNESSIONI", 12, GRAY, True, PP_ALIGN.CENTER, "Arial Black")

    add_textbox(sl, left_x, Inches(1.5), col_w, Inches(0.4),
                "🔴  1  AGENTE BASE", 14, RED, True)
    rect = add_rounded_rect(sl, left_x, Inches(2.0), col_w, Inches(3.8), LIGHT_GRAY)
    rect.line.color.rgb = RED
    add_textbox(sl, left_x + Inches(0.2), Inches(2.1), col_w - Inches(0.4), Inches(3.5),
                "Non risultano email non lette importanti ricevute oggi nei dati disponibili. "
                "Se hai bisogno di un riepilogo su un periodo diverso o su una categoria specifica di email, fammi sapere!",
                14, BLACK)
    add_textbox(sl, left_x, Inches(5.9), col_w, Inches(0.5),
                "🔴 Nessun accesso alla mail", 13, RED, True)

    add_textbox(sl, right_x, Inches(1.5), col_w, Inches(0.4),
                "🟡  2  CONNESSIONI", 14, YELLOW, True)
    rect = add_rounded_rect(sl, right_x, Inches(2.0), col_w, Inches(3.8), LIGHT_GRAY)
    rect.line.color.rgb = YELLOW
    add_textbox(sl, right_x + Inches(0.2), Inches(2.1), col_w - Inches(0.4), Inches(3.5),
                "1. PR: fix_validation_pipeline into main (Azure DevOps)\n"
                "2. Annullato: Sprint Review - Weekly Sync (Stefano Galli)\n"
                "3. [API Gateway] Integrazione Meridian — kick-off tecnico (Francesca Colombo)\n"
                "4. Analytics Dashboard [TEST] - Issue Report (errore critico)",
                13, BLACK)
    add_textbox(sl, right_x, Inches(5.9), col_w, Inches(0.5),
                "🟡 Elenca email categorizzate. Non distingue le priorità.", 12, YELLOW, True)

    # ── SLIDE 11: PROMPT 2 — Agent 3 vs 4 ──
    sl = prs.slides.add_slide(blank)
    add_textbox(sl, Inches(0), Inches(0.3), SLIDE_WIDTH, Inches(0.5),
                "PROMPT 2 — RISPOSTE A CONFRONTO", 22, BLACK, True, PP_ALIGN.CENTER, "Arial Black")
    add_textbox(sl, Inches(0), Inches(0.8), SLIDE_WIDTH, Inches(0.4),
                "WORK IQ COPILOT vs COMPLETO", 12, GRAY, True, PP_ALIGN.CENTER, "Arial Black")

    add_textbox(sl, left_x, Inches(1.5), col_w, Inches(0.4),
                "🟣  3  WORK IQ COPILOT", 14, PURPLE, True)
    rect = add_rounded_rect(sl, left_x, Inches(2.0), col_w, Inches(3.8), LIGHT_GRAY)
    rect.line.color.rgb = PURPLE
    add_textbox(sl, left_x + Inches(0.2), Inches(2.1), col_w - Inches(0.4), Inches(3.5),
                "Non risultano email non lette di oggi.\n\n"
                "Posso però:\n"
                "▪ Cercare nelle ultime 24-48 ore\n"
                "▪ Includere email già lette ma rilevanti\n"
                "▪ Filtrare per mittenti chiave (team Phoenix, Marco, Chiara)",
                13, BLACK)
    add_textbox(sl, left_x, Inches(5.9), col_w, Inches(0.5),
                "🟣 Spiega il contesto, propone strategie alternative intelligenti", 12, PURPLE, True)

    add_textbox(sl, right_x, Inches(1.5), col_w, Inches(0.4),
                "🟢  4  COMPLETO", 14, GREEN, True)
    rect = add_rounded_rect(sl, right_x, Inches(2.0), col_w, Inches(3.8), LIGHT_GRAY)
    rect.line.color.rgb = GREEN
    add_textbox(sl, right_x + Inches(0.2), Inches(2.1), col_w - Inches(0.4), Inches(3.5),
                "▪ Roberto Neri (29 mag) — Invito riunione Teams\n"
                "▪ Deskbooking (29 mag) — Check-in Reminder ×2\n"
                "▪ Microsoft Outlook (28 mag) — Non recapitabile: Follow-up con Sofia Marchetti ×2\n\n"
                "Vuoi che approfondisca per estrarre azioni?",
                13, BLACK)
    add_textbox(sl, right_x, Inches(5.9), col_w, Inches(0.5),
                "🟢 Email filtrate con data, anteprima e azioni di follow-up", 12, GREEN, True)

    # ── SLIDE 12: SECTION PROMPT 3 ──
    sl = prs.slides.add_slide(blank)
    add_textbox(sl, Inches(0), Inches(2.2), SLIDE_WIDTH, Inches(0.8),
                "PROMPT 3", 36, GREEN, True, PP_ALIGN.CENTER, "Arial Black")
    add_textbox(sl, Inches(1.5), Inches(3.5), Inches(10.3), Inches(1.5),
                '"Manda una email a s.marchetti@reply.it: digli che ho visto le sue domande sul progetto e ne parliamo alla prossima riunione."',
                18, GRAY, False, PP_ALIGN.CENTER)

    # ── SLIDE 13: PROMPT 3 — AZIONE ──
    sl = prs.slides.add_slide(blank)
    add_textbox(sl, Inches(0), Inches(0.3), SLIDE_WIDTH, Inches(0.5),
                "AGENTE 4 — AZIONE COMPLETATA", 22, BLACK, True, PP_ALIGN.CENTER, "Arial Black")
    add_textbox(sl, Inches(0), Inches(0.8), SLIDE_WIDTH, Inches(0.4),
                "CAPIRE + AGIRE = VALORE REALE", 12, GRAY, True, PP_ALIGN.CENTER, "Arial Black")

    card_x = Inches(3)
    card_w = Inches(7.3)
    rect = add_rounded_rect(sl, card_x, Inches(1.8), card_w, Inches(4.2), LIGHT_GRAY, GREEN)
    _, tf = add_rich_textbox(sl, card_x + Inches(0.4), Inches(2.0), card_w - Inches(0.8), Inches(3.8))
    tf.paragraphs[0].space_after = Pt(0)
    add_paragraph(tf, "📨 A: s.marchetti@reply.it", 16, GREEN, True, space_after=Pt(8))
    add_paragraph(tf, 'Oggetto: "Discussione delle domande sul progetto"', 13, GRAY, italic=True, space_after=Pt(16))
    add_paragraph(tf, "Ciao Sofia,", 14, BLACK, space_after=Pt(4))
    add_paragraph(tf, "ho visto le tue domande sul progetto.", 14, BLACK, space_after=Pt(4))
    add_paragraph(tf, "Ne parleremo nella prossima riunione.", 14, BLACK, space_after=Pt(16))
    add_paragraph(tf, "✅ Email inviata con successo", 16, GREEN, True, space_after=Pt(8))

    add_textbox(sl, Inches(2.5), Inches(6.2), Inches(8.3), Inches(0.8),
                "L'agente ha compreso il contesto dalla conversazione, formulato un messaggio appropriato e inviato l'email — tutto in un unico passaggio.",
                13, GRAY, False, PP_ALIGN.CENTER)

    # ── SLIDE 14: TABELLA COMPARATIVA ──
    sl = prs.slides.add_slide(blank)
    add_textbox(sl, Inches(0), Inches(0.3), SLIDE_WIDTH, Inches(0.5),
                "RIEPILOGO COMPARATIVO", 26, BLACK, True, PP_ALIGN.CENTER, "Arial Black")
    add_textbox(sl, Inches(0), Inches(0.85), SLIDE_WIDTH, Inches(0.4),
                "CAPIRE · AGIRE · FONTI", 12, GRAY, True, PP_ALIGN.CENTER, "Arial Black")

    # Build table
    rows, cols = 5, 5
    tbl_w = Inches(10)
    tbl_x = (SLIDE_WIDTH - tbl_w) // 2
    table_shape = sl.shapes.add_table(rows, cols, tbl_x, Inches(1.8), tbl_w, Inches(3.5))
    tbl = table_shape.table

    headers = ["#", "Livello", "Capisce", "Agisce", "Fonti dati"]
    for i, h in enumerate(headers):
        cell = tbl.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = GREEN
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.color.rgb = WHITE
                r.font.bold = True
                r.font.size = Pt(12)
                r.font.name = "Arial Black"

    data = [
        ["🔴", "Base", "❌", "❌", "Nessuna"],
        ["🟡", "Connessioni", "Parziale", "✅", "Mail, Calendar, User"],
        ["🟣", "Work IQ Copilot", "✅✅", "❌", "Semantic Index, Memory, M365"],
        ["🟢", "Completo", "✅✅", "✅", "Tutto"],
    ]
    for r_idx, row_data in enumerate(data):
        for c_idx, val in enumerate(row_data):
            cell = tbl.cell(r_idx + 1, c_idx)
            cell.text = val
            if r_idx % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(12)
                    r.font.name = "Arial"
                    if c_idx == 1:
                        r.font.bold = True

    # ── SLIDE 15: MESSAGGI CHIAVE ──
    sl = prs.slides.add_slide(blank)
    add_textbox(sl, Inches(0), Inches(0.3), SLIDE_WIDTH, Inches(0.5),
                "MESSAGGI CHIAVE", 26, BLACK, True, PP_ALIGN.CENTER, "Arial Black")

    messages = [
        ("Senza contesto", 'Un LLM generico non serve in azienda. È intelligente ma cieco.', '"Non risultano informazioni…"'),
        ("Connessioni", 'Con le API vede le tue app e fa cose. Ma cerca per keyword, non capisce il contesto.', '"Ecco i dettagli della tua riunione"'),
        ("Semantic Index", 'Il semantic index connette i puntini nel tempo. Analisi strutturata e proattiva.', '"Preparati su rischi e dipendenze"'),
        ("Completo", 'Capire + Agire = valore reale. Ha inviato l\'email.', '"Email inviata con successo ✅"'),
        ("Governance", 'Ogni tool è controllabile dall\'admin center. Ogni chiamata è tracciabile.', None),
        ("MCP Standard", 'Domani aggiungi SharePoint, Dynamics, i tuoi sistemi custom. Protocollo aperto.', None),
    ]
    msg_w = Inches(5.8)
    msg_h = Inches(1.8)
    gap_x = Inches(0.4)
    gap_y = Inches(0.2)
    start_x = (SLIDE_WIDTH - msg_w * 2 - gap_x) // 2
    start_y = Inches(1.3)
    for i, (label, text, quote) in enumerate(messages):
        col = i % 2
        row = i // 2
        x = start_x + col * (msg_w + gap_x)
        y = start_y + row * (msg_h + gap_y)
        rect = add_rounded_rect(sl, x, y, msg_w, msg_h, LIGHT_GRAY)
        _, tf = add_rich_textbox(sl, x + Inches(0.3), y + Inches(0.2), msg_w - Inches(0.6), msg_h - Inches(0.3))
        tf.paragraphs[0].space_after = Pt(0)
        add_paragraph(tf, label, 10, GREEN, True, font_name="Arial Black", space_after=Pt(4))
        add_paragraph(tf, text, 12, BLACK, space_after=Pt(4))
        if quote:
            add_paragraph(tf, quote, 11, GRAY, italic=True, space_after=Pt(0))

    # ── SLIDE 16: THANK YOU ──
    sl = prs.slides.add_slide(blank)
    set_slide_gradient(sl, DARK_GREEN, GREEN)
    add_textbox(sl, Inches(0), Inches(2.5), SLIDE_WIDTH, Inches(1.5),
                "THANK YOU", 60, WHITE, True, PP_ALIGN.CENTER, "Arial Black")
    add_textbox(sl, Inches(0), Inches(4.5), SLIDE_WIDTH, Inches(0.5),
                "www.reply.com", 16, RGBColor(0xCC, 0xCC, 0xCC), False, PP_ALIGN.CENTER)

    out_path = os.path.join(SCRIPT_DIR, "presentazione.pptx")
    prs.save(out_path)
    print(f"✅ Salvato: {out_path}")
    return out_path


# ═══════════════════════════════════════════════════════
#  VISUAL PRESENTATION  (presentazione_visual.pptx)
# ═══════════════════════════════════════════════════════

def build_visual_pptx():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank = prs.slide_layouts[6]

    # ── SLIDE 1: COVER ──
    sl = prs.slides.add_slide(blank)
    set_slide_gradient(sl, DARK_GREEN, GREEN)
    add_textbox(sl, Inches(1), Inches(3.5), Inches(9), Inches(2),
                "WORK IQ\nDEMO\nCOMPARATIVA", 44, WHITE, True, font_name="Arial Black")
    add_textbox(sl, Inches(1), Inches(5.5), Inches(9), Inches(0.5),
                "4 livelli di intelligenza — Stesso prompt, risposte diverse", 16, WHITE)
    add_textbox(sl, Inches(1), Inches(6.1), Inches(9), Inches(0.4),
                "1 Giugno 2026 · Copilot Studio + Work IQ MCP", 13, RGBColor(0xCC, 0xCC, 0xCC))

    # ── SLIDE 2: STATEMENT ──
    sl = prs.slides.add_slide(blank)
    set_slide_bg(sl, BLACK)
    add_textbox(sl, Inches(1), Inches(2), Inches(11), Inches(3.5),
                "SENZA CONTESTO,\nL'AI NON SERVE\nIN AZIENDA.", 40, WHITE, True, PP_ALIGN.CENTER, "Arial Black")
    add_textbox(sl, Inches(4.5), Inches(3.35), Inches(4.5), Inches(0.8),
                "NON SERVE", 40, GREEN, True, PP_ALIGN.CENTER, "Arial Black")

    # ── SLIDE 3: 4 LIVELLI ──
    sl = prs.slides.add_slide(blank)
    add_textbox(sl, Inches(0), Inches(0.4), SLIDE_WIDTH, Inches(0.6),
                "I 4 LIVELLI DI INTELLIGENZA", 26, BLACK, True, PP_ALIGN.CENTER, "Arial Black")
    add_textbox(sl, Inches(0), Inches(1.0), SLIDE_WIDTH, Inches(0.4),
                "STESSO PROMPT, CAPACITÀ CRESCENTI", 13, GRAY, True, PP_ALIGN.CENTER, "Arial Black")
    cards = [
        ("1", "BASE", "Solo LLM generico.\nNessun tool, nessun dato.", "—", RED),
        ("2", "CONNESSIONI", "Accesso diretto a dati\nMicrosoft 365.", "Mail · Calendar · User", YELLOW),
        ("3", "WORK IQ COPILOT", "Semantic Index: capisce\nil contesto nel tempo.", "Semantic Index · Memory\nCross-app", PURPLE),
        ("4", "COMPLETO", "Intelligenza profonda +\naccesso diretto.\nCapisce e agisce.", "Copilot + Mail +\nCalendar + User", GREEN),
    ]
    card_w = Inches(2.8)
    gap = Inches(0.3)
    total = card_w * 4 + gap * 3
    start_x = (SLIDE_WIDTH - total) // 2
    for i, (num, name, desc, tools, color) in enumerate(cards):
        x = start_x + i * (card_w + gap)
        y = Inches(1.8)
        h = Inches(4.5)
        add_rounded_rect(sl, x, y, card_w, h, LIGHT_GRAY, color)
        bar = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, card_w, Inches(0.15))
        bar.fill.solid()
        bar.fill.fore_color.rgb = color
        bar.line.fill.background()
        add_textbox(sl, x, y + Inches(0.4), card_w, Inches(0.7),
                    num, 32, GRAY, True, PP_ALIGN.CENTER, "Arial Black")
        add_textbox(sl, x, y + Inches(1.2), card_w, Inches(0.5),
                    name, 14, BLACK, True, PP_ALIGN.CENTER, "Arial Black")
        add_textbox(sl, x + Inches(0.2), y + Inches(1.9), card_w - Inches(0.4), Inches(1.5),
                    desc, 12, GRAY, False, PP_ALIGN.CENTER)
        add_textbox(sl, x + Inches(0.2), y + Inches(3.5), card_w - Inches(0.4), Inches(0.8),
                    tools, 10, GRAY, False, PP_ALIGN.CENTER)

    # ── SCREENSHOT SLIDE HELPER ──
    def screenshot_slide(title, subtitle, img_file):
        """Full-width screenshot on dark background."""
        sl = prs.slides.add_slide(blank)
        set_slide_bg(sl, BG_DARK)
        add_textbox(sl, Inches(0), Inches(0.3), SLIDE_WIDTH, Inches(0.5),
                    title, 22, WHITE, True, PP_ALIGN.CENTER, "Arial Black")
        add_textbox(sl, Inches(0), Inches(0.85), SLIDE_WIDTH, Inches(0.35),
                    subtitle, 11, RGBColor(0x8B, 0x94, 0x9E), True, PP_ALIGN.CENTER, "Arial Black")
        img_path = os.path.join(SCREENSHOTS_DIR, img_file)
        if os.path.exists(img_path):
            add_image_centered(sl, img_path,
                               Inches(12.5), Inches(5.8), Inches(1.4))
        return sl

    def duo_slide(title, subtitle, img1, label1, badge1, img2, label2, badge2):
        """Two screenshots side by side on dark background."""
        sl = prs.slides.add_slide(blank)
        set_slide_bg(sl, BG_DARK)
        add_textbox(sl, Inches(0), Inches(0.3), SLIDE_WIDTH, Inches(0.5),
                    title, 20, WHITE, True, PP_ALIGN.CENTER, "Arial Black")
        add_textbox(sl, Inches(0), Inches(0.8), SLIDE_WIDTH, Inches(0.35),
                    subtitle, 11, RGBColor(0x8B, 0x94, 0x9E), True, PP_ALIGN.CENTER, "Arial Black")

        img_w = Inches(6)
        img_h = Inches(5)
        gap = Inches(0.3)
        left_x = (SLIDE_WIDTH - img_w * 2 - gap) // 2
        right_x = left_x + img_w + gap
        top_y = Inches(1.4)

        for img_file, label, badge_text, badge_color, x in [
            (img1, label1, badge1[0], badge1[1], left_x),
            (img2, label2, badge2[0], badge2[1], right_x),
        ]:
            img_path = os.path.join(SCREENSHOTS_DIR, img_file)
            if os.path.exists(img_path):
                from PIL import Image as PILImage
                with PILImage.open(img_path) as pimg:
                    pw, ph = pimg.size
                aspect = pw / ph
                actual_w = img_w
                actual_h = int(actual_w / aspect)
                if actual_h > img_h:
                    actual_h = img_h
                    actual_w = int(actual_h * aspect)
                sl.shapes.add_picture(img_path, x, top_y, actual_w, actual_h)

            # Badge
            badge_rect = add_rounded_rect(sl, x, top_y + img_h + Inches(0.15),
                                          Inches(2.2), Inches(0.35), badge_color)
            badge_rect.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            run = badge_rect.text_frame.paragraphs[0].add_run()
            run.text = badge_text
            run.font.size = Pt(9)
            run.font.color.rgb = WHITE
            run.font.bold = True
            run.font.name = "Arial Black"

            # Label
            add_textbox(sl, x, top_y + img_h + Inches(0.55),
                        img_w, Inches(0.4), label, 11, RGBColor(0xCC, 0xCC, 0xCC))

        return sl

    # ── SLIDE 4: DASHBOARD ──
    screenshot_slide("DASHBOARD DI PARTENZA",
                     "4 AGENTI COPILOT STUDIO CONNESSI IN PARALLELO",
                     "00_dashboard_pulita.png")

    # ── SLIDE 5: SECTION PROMPT 1 ──
    sl = prs.slides.add_slide(blank)
    add_textbox(sl, Inches(0), Inches(2.5), SLIDE_WIDTH, Inches(0.8),
                "PROMPT 1", 36, GREEN, True, PP_ALIGN.CENTER, "Arial Black")
    add_textbox(sl, Inches(2), Inches(3.8), Inches(9.3), Inches(1),
                '"Cosa devo sapere prima della mia prossima riunione?"',
                20, GRAY, False, PP_ALIGN.CENTER)

    # ── SLIDE 6: P1 OVERVIEW ──
    screenshot_slide("PROMPT 1 — VISTA COMPARATIVA",
                     "STESSO PROMPT, 4 RISPOSTE DIVERSE",
                     "01_prompt1_overview.png")

    # ── SLIDE 7: P1 AGENT 1 vs 2 ──
    duo_slide("PROMPT 1 — BASE vs CONNESSIONI",
              'DA "NON HO INFORMAZIONI" A DETTAGLI COMPLETI',
              "02_prompt1_agent1.png", "🔴 Educato ma inutile. Non sa nulla del tuo lavoro.",
              ("1 BASE", RED),
              "02_prompt1_agent2.png", "🟡 Vede il calendario: titolo, orario, partecipanti.",
              ("2 CONNESSIONI", YELLOW))

    # ── SLIDE 8: P1 AGENT 3 vs 4 ──
    duo_slide("PROMPT 1 — COPILOT vs COMPLETO",
              "ANALISI PROFONDA E SUGGERIMENTI OPERATIVI",
              "02_prompt1_agent3.png", "🟣 Analisi strutturata, contesto storico, briefing operativo.",
              ("3 WORK IQ COPILOT", PURPLE),
              "02_prompt1_agent4.png", "🟢 Dettagli operativi + stato RSVP + azioni concrete.",
              ("4 COMPLETO", GREEN))

    # ── SLIDE 9: SECTION PROMPT 2 ──
    sl = prs.slides.add_slide(blank)
    add_textbox(sl, Inches(0), Inches(2.5), SLIDE_WIDTH, Inches(0.8),
                "PROMPT 2", 36, GREEN, True, PP_ALIGN.CENTER, "Arial Black")
    add_textbox(sl, Inches(2), Inches(3.8), Inches(9.3), Inches(1),
                '"Riassumi le email non lette più importanti di oggi."',
                20, GRAY, False, PP_ALIGN.CENTER)

    # ── SLIDE 10: P2 OVERVIEW ──
    screenshot_slide("PROMPT 2 — VISTA COMPARATIVA",
                     "EMAIL: DA ZERO INFORMAZIONI A INSIGHT STRUTTURATI",
                     "06_prompt2_overview.png")

    # ── SLIDE 11: P2 AGENT 1 vs 2 ──
    duo_slide("PROMPT 2 — BASE vs CONNESSIONI",
              "ACCESSO AI DATI VS ANALISI CIECA",
              "07_prompt2_agent1.png", "🔴 Nessun accesso alla mail.",
              ("1 BASE", RED),
              "07_prompt2_agent2.png", "🟡 Elenca email categorizzate con azioni suggerite.",
              ("2 CONNESSIONI", YELLOW))

    # ── SLIDE 12: P2 AGENT 3 vs 4 ──
    duo_slide("PROMPT 2 — COPILOT vs COMPLETO",
              "STRATEGIE ALTERNATIVE VS FILTRAGGIO INTELLIGENTE",
              "07_prompt2_agent3.png", "🟣 Spiega l'assenza di risultati e propone strategie alternative.",
              ("3 WORK IQ COPILOT", PURPLE),
              "07_prompt2_agent4.png", "🟢 Email filtrate per rilevanza con data e anteprima.",
              ("4 COMPLETO", GREEN))

    # ── SLIDE 13: SECTION PROMPT 3 ──
    sl = prs.slides.add_slide(blank)
    add_textbox(sl, Inches(0), Inches(2.2), SLIDE_WIDTH, Inches(0.8),
                "PROMPT 3", 36, GREEN, True, PP_ALIGN.CENTER, "Arial Black")
    add_textbox(sl, Inches(1.5), Inches(3.5), Inches(10.3), Inches(1.5),
                '"Manda una email a s.marchetti@reply.it: digli che ho visto le sue domande sul progetto e ne parliamo alla prossima riunione."',
                18, GRAY, False, PP_ALIGN.CENTER)

    # ── SLIDE 14: P3 AGENT 4 EMAIL — Annotated ──
    sl = prs.slides.add_slide(blank)
    set_slide_bg(sl, BG_DARK)
    add_textbox(sl, Inches(0), Inches(0.3), SLIDE_WIDTH, Inches(0.5),
                "AGENTE 4 — AZIONE COMPLETATA", 22, WHITE, True, PP_ALIGN.CENTER, "Arial Black")
    add_textbox(sl, Inches(0), Inches(0.8), SLIDE_WIDTH, Inches(0.35),
                "CAPIRE + AGIRE = VALORE REALE", 11, RGBColor(0x8B, 0x94, 0x9E), True, PP_ALIGN.CENTER, "Arial Black")

    # Screenshot on left (2/3)
    img_path = os.path.join(SCREENSHOTS_DIR, "12_prompt3_agent4_email.png")
    if os.path.exists(img_path):
        sl.shapes.add_picture(img_path, Inches(0.5), Inches(1.5), Inches(8), Inches(5.5))

    # Annotation on right (1/3)
    _, tf = add_rich_textbox(sl, Inches(9), Inches(1.8), Inches(3.8), Inches(5))
    tf.paragraphs[0].space_after = Pt(0)
    add_paragraph(tf, "UN SINGOLO PROMPT.", 16, GREEN, True, font_name="Arial Black", space_after=Pt(12))
    add_paragraph(tf, "L'agente ha compreso il contesto dalla conversazione, formulato un messaggio appropriato e inviato l'email.",
                  13, RGBColor(0x8B, 0x94, 0x9E), space_after=Pt(20))
    p = add_paragraph(tf, "", space_after=Pt(4))
    add_run(p, "1", 48, WHITE, True, font_name="Arial Black")
    add_paragraph(tf, "PASSAGGIO PER CAPIRE E AGIRE", 10, RGBColor(0x8B, 0x94, 0x9E), True, font_name="Arial Black", space_after=Pt(24))
    add_paragraph(tf, "📨 A: s.marchetti@reply.it", 12, WHITE, space_after=Pt(6))
    add_paragraph(tf, "✅ Email inviata con successo", 13, GREEN, True, space_after=Pt(0))

    # ── SLIDE 15: PANORAMICA FINALE ──
    screenshot_slide("PANORAMICA FINALE",
                     "TUTTA LA CONVERSAZIONE A CONFRONTO",
                     "13_panoramica_finale.png")

    # ── SLIDE 16: TABELLA COMPARATIVA ──
    sl = prs.slides.add_slide(blank)
    add_textbox(sl, Inches(0), Inches(0.3), SLIDE_WIDTH, Inches(0.5),
                "RIEPILOGO COMPARATIVO", 26, BLACK, True, PP_ALIGN.CENTER, "Arial Black")
    add_textbox(sl, Inches(0), Inches(0.85), SLIDE_WIDTH, Inches(0.4),
                "CAPIRE · AGIRE · FONTI", 12, GRAY, True, PP_ALIGN.CENTER, "Arial Black")
    rows, cols = 5, 5
    tbl_w = Inches(10)
    tbl_x = (SLIDE_WIDTH - tbl_w) // 2
    table_shape = sl.shapes.add_table(rows, cols, tbl_x, Inches(1.8), tbl_w, Inches(3.5))
    tbl = table_shape.table
    headers = ["#", "Livello", "Capisce", "Agisce", "Fonti dati"]
    for i, h in enumerate(headers):
        cell = tbl.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = GREEN
        for p in cell.text_frame.paragraphs:
            for r in p.runs:
                r.font.color.rgb = WHITE
                r.font.bold = True
                r.font.size = Pt(12)
                r.font.name = "Arial Black"
    data = [
        ["🔴", "Base", "❌", "❌", "Nessuna"],
        ["🟡", "Connessioni", "Parziale", "✅", "Mail, Calendar, User"],
        ["🟣", "Work IQ Copilot", "✅✅", "❌", "Semantic Index, Memory, M365"],
        ["🟢", "Completo", "✅✅", "✅", "Tutto"],
    ]
    for r_idx, row_data in enumerate(data):
        for c_idx, val in enumerate(row_data):
            cell = tbl.cell(r_idx + 1, c_idx)
            cell.text = val
            if r_idx % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY
            for p in cell.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(12)
                    r.font.name = "Arial"
                    if c_idx == 1:
                        r.font.bold = True

    # ── SLIDE 17: MESSAGGI CHIAVE ──
    sl = prs.slides.add_slide(blank)
    add_textbox(sl, Inches(0), Inches(0.3), SLIDE_WIDTH, Inches(0.5),
                "MESSAGGI CHIAVE", 26, BLACK, True, PP_ALIGN.CENTER, "Arial Black")
    messages = [
        ("Senza contesto", 'Un LLM generico non serve in azienda. È intelligente ma cieco.', '"Non risultano informazioni…"'),
        ("Connessioni", 'Con le API vede le tue app e fa cose. Ma cerca per keyword, non capisce il contesto.', '"Ecco i dettagli della tua riunione"'),
        ("Semantic Index", 'Il semantic index connette i puntini nel tempo. Analisi strutturata e proattiva.', '"Preparati su rischi e dipendenze"'),
        ("Completo", 'Capire + Agire = valore reale. Ha inviato l\'email.', '"Email inviata con successo ✅"'),
        ("Governance", 'Ogni tool è controllabile dall\'admin center. Ogni chiamata è tracciabile.', None),
        ("MCP Standard", 'Domani aggiungi SharePoint, Dynamics, i tuoi sistemi custom. Protocollo aperto.', None),
    ]
    msg_w = Inches(5.8)
    msg_h = Inches(1.8)
    gap_x = Inches(0.4)
    gap_y = Inches(0.2)
    start_x = (SLIDE_WIDTH - msg_w * 2 - gap_x) // 2
    start_y = Inches(1.3)
    for i, (label, text, quote) in enumerate(messages):
        col = i % 2
        row = i // 2
        x = start_x + col * (msg_w + gap_x)
        y = start_y + row * (msg_h + gap_y)
        add_rounded_rect(sl, x, y, msg_w, msg_h, LIGHT_GRAY)
        _, tf = add_rich_textbox(sl, x + Inches(0.3), y + Inches(0.2), msg_w - Inches(0.6), msg_h - Inches(0.3))
        tf.paragraphs[0].space_after = Pt(0)
        add_paragraph(tf, label, 10, GREEN, True, font_name="Arial Black", space_after=Pt(4))
        add_paragraph(tf, text, 12, BLACK, space_after=Pt(4))
        if quote:
            add_paragraph(tf, quote, 11, GRAY, italic=True, space_after=Pt(0))

    # ── SLIDE 18: THANK YOU ──
    sl = prs.slides.add_slide(blank)
    set_slide_gradient(sl, DARK_GREEN, GREEN)
    add_textbox(sl, Inches(0), Inches(2.5), SLIDE_WIDTH, Inches(1.5),
                "THANK YOU", 60, WHITE, True, PP_ALIGN.CENTER, "Arial Black")
    add_textbox(sl, Inches(0), Inches(4.5), SLIDE_WIDTH, Inches(0.5),
                "www.reply.com", 16, RGBColor(0xCC, 0xCC, 0xCC), False, PP_ALIGN.CENTER)

    out_path = os.path.join(SCRIPT_DIR, "presentazione_visual.pptx")
    prs.save(out_path)
    print(f"✅ Salvato: {out_path}")
    return out_path


# ═══════════════════════════════════════════════════════
#  MAIN
# ═══════════════════════════════════════════════════════
if __name__ == "__main__":
    print("Generazione PPTX in corso...\n")
    build_text_pptx()
    build_visual_pptx()
    print("\n🎉 Entrambi i file PPTX sono stati generati!")
