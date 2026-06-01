#!/usr/bin/env python3
"""
Genera due file PPTX dalla presentazione Work IQ Demo Comparativa:
  1. presentazione.pptx          — versione testuale (senza immagini)
  2. presentazione_visual.pptx   — versione con screenshot integrati

Richiede: pip install python-pptx Pillow
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Paths ──
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
SCREENSHOTS_DIR = os.path.join(SCRIPT_DIR, "screenshots")

# ── Brand Colors ──
GREEN       = RGBColor(0x76, 0xB8, 0x2A)
DARK_GREEN  = RGBColor(0x2D, 0x6A, 0x2E)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
BLACK       = RGBColor(0x1A, 0x1A, 0x1A)
GRAY        = RGBColor(0x55, 0x55, 0x55)
LIGHT_GRAY  = RGBColor(0xF5, 0xF5, 0xF5)
MID_GRAY    = RGBColor(0xE8, 0xE8, 0xE8)
BG_DARK     = RGBColor(0x0D, 0x11, 0x17)
RED         = RGBColor(0xEF, 0x44, 0x44)
YELLOW      = RGBColor(0xF5, 0x9E, 0x0B)
PURPLE      = RGBColor(0x8B, 0x5C, 0xF6)
BRIGHT_GREEN = RGBColor(0x8C, 0xC6, 0x3F)
MUTED_TEXT  = RGBColor(0x8B, 0x94, 0x9E)

# ── Standard slide dimensions (16:9 widescreen) ──
SLIDE_WIDTH  = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# ── Layout constants ──
MARGIN_H = Inches(0.8)          # horizontal margin
CONTENT_W = SLIDE_WIDTH - 2 * MARGIN_H  # usable width
TITLE_Y = Inches(0.5)           # title vertical pos
SUBTITLE_Y = Inches(1.1)        # subtitle vertical pos
BODY_TOP = Inches(1.7)          # content area starts


# ═══════════════════════════════════════════════════════
# Utility helpers
# ═══════════════════════════════════════════════════════

def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def set_slide_gradient(slide, color1, color2):
    bg = slide.background
    fill = bg.fill
    fill.gradient()
    fill.gradient_stops[0].color.rgb = color1
    fill.gradient_stops[0].position = 0.0
    fill.gradient_stops[1].color.rgb = color2
    fill.gradient_stops[1].position = 1.0


def add_textbox(slide, left, top, width, height, text, font_size=18,
                color=BLACK, bold=False, alignment=PP_ALIGN.LEFT,
                font_name="Arial", anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font_name
    return txBox, tf


def add_rich_textbox(slide, left, top, width, height, anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    return txBox, tf


def add_run(paragraph, text, font_size=14, color=BLACK, bold=False,
            italic=False, font_name="Arial"):
    run = paragraph.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font_name
    return run


def add_paragraph(tf, text="", font_size=14, color=BLACK, bold=False,
                  italic=False, alignment=PP_ALIGN.LEFT, font_name="Arial",
                  space_after=Pt(6), space_before=Pt(0)):
    p = tf.add_paragraph()
    p.alignment = alignment
    p.space_after = space_after
    p.space_before = space_before
    if text:
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.italic = italic
        run.font.name = font_name
    return p


def add_rounded_rect(slide, left, top, width, height, fill_color, border_color=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(2)
    else:
        shape.line.fill.background()
    return shape


def add_color_bar(slide, left, top, width, color, height=Inches(0.12)):
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = color
    bar.line.fill.background()
    return bar


def add_image_fitted(slide, img_path, max_width, max_height, top_offset, center_h=True):
    """Add image preserving aspect ratio, centered horizontally."""
    from PIL import Image as PILImage
    with PILImage.open(img_path) as img:
        img_w, img_h = img.size
    aspect = img_w / img_h
    target_w = max_width
    target_h = int(target_w / aspect)
    if target_h > max_height:
        target_h = max_height
        target_w = int(target_h * aspect)
    left = (SLIDE_WIDTH - target_w) // 2 if center_h else MARGIN_H
    return slide.shapes.add_picture(img_path, left, top_offset, target_w, target_h)


def add_slide_title(slide, title, subtitle=None, color=BLACK, sub_color=GRAY):
    """Standard title + subtitle block at top of slide."""
    add_textbox(slide, MARGIN_H, TITLE_Y, CONTENT_W, Inches(0.6),
                title, 24, color, True, PP_ALIGN.LEFT, "Arial Black")
    if subtitle:
        add_textbox(slide, MARGIN_H, SUBTITLE_Y, CONTENT_W, Inches(0.35),
                    subtitle, 12, sub_color, True, PP_ALIGN.LEFT, "Arial Black")


def add_slide_title_center(slide, title, subtitle=None, color=BLACK, sub_color=GRAY):
    add_textbox(slide, Inches(0), TITLE_Y, SLIDE_WIDTH, Inches(0.6),
                title, 24, color, True, PP_ALIGN.CENTER, "Arial Black")
    if subtitle:
        add_textbox(slide, Inches(0), SUBTITLE_Y, SLIDE_WIDTH, Inches(0.35),
                    subtitle, 12, sub_color, True, PP_ALIGN.CENTER, "Arial Black")


# ═══════════════════════════════════════════════════════
#  SHARED BUILDERS — reused across both decks
# ═══════════════════════════════════════════════════════

def build_cover(sl):
    set_slide_gradient(sl, DARK_GREEN, GREEN)
    add_textbox(sl, Inches(1.5), Inches(2.2), Inches(8), Inches(2.8),
                "WORK IQ\nDEMO\nCOMPARATIVA", 52, WHITE, True, font_name="Arial Black")
    add_textbox(sl, Inches(1.5), Inches(5.2), Inches(8), Inches(0.5),
                "4 livelli di intelligenza — Stesso prompt, risposte diverse", 18, WHITE)
    add_textbox(sl, Inches(1.5), Inches(5.9), Inches(8), Inches(0.4),
                "1 Giugno 2026 · Copilot Studio + Work IQ MCP", 14, RGBColor(0xBB, 0xDD, 0xBB))


def build_statement(sl):
    set_slide_bg(sl, BLACK)
    # All text in single textbox with colored "NON SERVE"
    _, tf = add_rich_textbox(sl, Inches(1.5), Inches(2.2), Inches(10.3), Inches(3.5))
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    tf.paragraphs[0].space_after = Pt(0)
    add_run(tf.paragraphs[0], "SENZA CONTESTO,", 42, WHITE, True, font_name="Arial Black")
    p2 = add_paragraph(tf, alignment=PP_ALIGN.CENTER, space_after=Pt(0))
    add_run(p2, "L'AI ", 42, WHITE, True, font_name="Arial Black")
    add_run(p2, "NON SERVE", 42, GREEN, True, font_name="Arial Black")
    p3 = add_paragraph(tf, alignment=PP_ALIGN.CENTER, space_after=Pt(0))
    add_run(p3, "IN AZIENDA.", 42, WHITE, True, font_name="Arial Black")


def build_four_pillars(sl):
    add_slide_title_center(sl, "I 4 LIVELLI DI INTELLIGENZA", "STESSO PROMPT, CAPACITÀ CRESCENTI")
    cards = [
        ("1", "BASE", "Solo LLM generico.\nNessun tool, nessun dato.", "—", RED),
        ("2", "CONNESSIONI", "Accesso diretto a dati\nMicrosoft 365.", "Mail · Calendar · User", YELLOW),
        ("3", "WORK IQ COPILOT", "Semantic Index: capisce\nil contesto nel tempo.", "Semantic Index · Memory\nCross-app", PURPLE),
        ("4", "COMPLETO", "Intelligenza profonda +\naccesso diretto. Capisce e agisce.", "Copilot + Mail + Calendar + User", GREEN),
    ]
    card_w = Inches(2.7)
    gap = Inches(0.35)
    total = card_w * 4 + gap * 3
    start_x = (SLIDE_WIDTH - total) // 2
    card_top = Inches(1.8)
    card_h = Inches(5.0)
    for i, (num, name, desc, tools, color) in enumerate(cards):
        x = start_x + i * (card_w + gap)
        add_rounded_rect(sl, x, card_top, card_w, card_h, LIGHT_GRAY)
        add_color_bar(sl, x, card_top, card_w, color)
        # Number — large, faded
        add_textbox(sl, x, card_top + Inches(0.5), card_w, Inches(0.8),
                    num, 36, MID_GRAY, True, PP_ALIGN.CENTER, "Arial Black")
        # Name
        add_textbox(sl, x + Inches(0.15), card_top + Inches(1.4), card_w - Inches(0.3), Inches(0.5),
                    name, 13, BLACK, True, PP_ALIGN.CENTER, "Arial Black")
        # Description
        add_textbox(sl, x + Inches(0.25), card_top + Inches(2.1), card_w - Inches(0.5), Inches(1.6),
                    desc, 11, GRAY, False, PP_ALIGN.CENTER)
        # Tools footer
        add_textbox(sl, x + Inches(0.25), card_top + Inches(3.9), card_w - Inches(0.5), Inches(0.7),
                    tools, 9, GRAY, False, PP_ALIGN.CENTER)


def build_section_slide(sl, title, subtitle):
    add_textbox(sl, Inches(0), Inches(2.8), SLIDE_WIDTH, Inches(0.9),
                title, 40, GREEN, True, PP_ALIGN.CENTER, "Arial Black")
    add_textbox(sl, Inches(2), Inches(4.0), Inches(9.3), Inches(1.5),
                subtitle, 20, GRAY, False, PP_ALIGN.CENTER)


def build_comparison_slide(sl, title, subtitle, left_data, right_data):
    """
    left_data / right_data: (badge_emoji, badge_num, badge_name, badge_color, text, verdict_text, verdict_color)
    """
    add_slide_title_center(sl, title, subtitle)
    col_w = Inches(5.6)
    gap = Inches(0.5)
    left_x = (SLIDE_WIDTH - col_w * 2 - gap) // 2
    right_x = left_x + col_w + gap
    label_top = Inches(1.7)
    box_top = Inches(2.2)
    box_h = Inches(4.2)
    inner_pad = Inches(0.35)

    for x, (emoji, num, name, badge_color, text, verdict, verdict_color) in [
        (left_x, left_data), (right_x, right_data)
    ]:
        # Agent label row
        _, tf = add_rich_textbox(sl, x, label_top, col_w, Inches(0.45))
        tf.paragraphs[0].space_after = Pt(0)
        add_run(tf.paragraphs[0], f"{emoji}  {num}  {name}", 16, badge_color, True)

        # Response box with left accent border
        add_rounded_rect(sl, x, box_top, col_w, box_h, LIGHT_GRAY, badge_color)
        # Override left border: colored accent bar
        add_color_bar(sl, x, box_top, Inches(0.08), badge_color, box_h)

        # Response text — padded inside the box
        add_textbox(sl, x + inner_pad, box_top + Inches(0.3),
                    col_w - inner_pad * 2, box_h - Inches(0.6),
                    text, 14, BLACK)

        # Verdict below box
        add_textbox(sl, x, box_top + box_h + Inches(0.12), col_w, Inches(0.4),
                    verdict, 12, verdict_color, True)


def build_table_slide(sl):
    add_slide_title_center(sl, "RIEPILOGO COMPARATIVO", "CAPIRE · AGIRE · FONTI")

    rows, cols = 5, 5
    tbl_w = Inches(10.5)
    tbl_x = (SLIDE_WIDTH - tbl_w) // 2
    tbl_y = Inches(2.5)
    tbl_h = Inches(3.2)
    table_shape = sl.shapes.add_table(rows, cols, tbl_x, tbl_y, tbl_w, tbl_h)
    tbl = table_shape.table

    # Column widths: proportional
    col_widths = [Inches(0.8), Inches(2.5), Inches(1.8), Inches(1.8), Inches(3.6)]
    for i, w in enumerate(col_widths):
        tbl.columns[i].width = w

    headers = ["#", "Livello", "Capisce", "Agisce", "Fonti dati"]
    for i, h in enumerate(headers):
        cell = tbl.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = GREEN
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER if i != 4 else PP_ALIGN.LEFT
            for r in p.runs:
                r.font.color.rgb = WHITE
                r.font.bold = True
                r.font.size = Pt(13)
                r.font.name = "Arial Black"

    data = [
        ("🔴", "Base",             "❌",       "❌", "Nessuna"),
        ("🟡", "Connessioni",      "Parziale", "✅", "Mail, Calendar, User"),
        ("🟣", "Work IQ Copilot",  "✅✅",     "❌", "Semantic Index, Memory, M365"),
        ("🟢", "Completo",         "✅✅",     "✅", "Tutto"),
    ]
    for r_idx, row_data in enumerate(data):
        for c_idx, val in enumerate(row_data):
            cell = tbl.cell(r_idx + 1, c_idx)
            cell.text = val
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if r_idx % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if c_idx != 4 else PP_ALIGN.LEFT
                for r in p.runs:
                    r.font.size = Pt(13)
                    r.font.name = "Arial"
                    if c_idx == 1:
                        r.font.bold = True


def build_messages_slide(sl):
    add_slide_title_center(sl, "MESSAGGI CHIAVE", " ")  # thin subtitle for spacing
    messages = [
        ("Senza contesto", 'Un LLM generico non serve in azienda. È intelligente ma cieco.',
         '"Non risultano informazioni…"'),
        ("Connessioni", 'Con le API vede le tue app e fa cose. Ma cerca per keyword, non capisce il contesto.',
         '"Ecco i dettagli della tua riunione"'),
        ("Semantic Index", 'Il semantic index connette i puntini nel tempo. Analisi strutturata e proattiva.',
         '"Preparati su rischi e dipendenze"'),
        ("Completo", 'Capire + Agire = valore reale. Ha inviato l\'email.',
         '"Email inviata con successo ✅"'),
        ("Governance", 'Ogni tool è controllabile dall\'admin center. Ogni chiamata è tracciabile.', None),
        ("MCP Standard", 'Domani aggiungi SharePoint, Dynamics, i tuoi sistemi custom. Protocollo aperto.', None),
    ]
    msg_w = Inches(5.5)
    msg_h = Inches(1.65)
    gap_x = Inches(0.5)
    gap_y = Inches(0.25)
    grid_w = msg_w * 2 + gap_x
    start_x = (SLIDE_WIDTH - grid_w) // 2
    grid_h = msg_h * 3 + gap_y * 2
    start_y = (SLIDE_HEIGHT - grid_h) // 2 + Inches(0.3)  # vertically centered, offset for title
    for i, (label, text, quote) in enumerate(messages):
        col = i % 2
        row = i // 2
        x = start_x + col * (msg_w + gap_x)
        y = start_y + row * (msg_h + gap_y)
        add_rounded_rect(sl, x, y, msg_w, msg_h, LIGHT_GRAY)
        pad = Inches(0.35)
        _, tf = add_rich_textbox(sl, x + pad, y + Inches(0.2), msg_w - pad * 2, msg_h - Inches(0.3))
        tf.paragraphs[0].space_after = Pt(0)
        add_paragraph(tf, label, 11, GREEN, True, font_name="Arial Black", space_after=Pt(6))
        add_paragraph(tf, text, 12, BLACK, space_after=Pt(4))
        if quote:
            add_paragraph(tf, quote, 11, GRAY, italic=True, space_after=Pt(0))


def build_thankyou(sl, url="www.reply.com"):
    set_slide_gradient(sl, DARK_GREEN, GREEN)
    add_textbox(sl, Inches(0), Inches(2.7), SLIDE_WIDTH, Inches(1.5),
                "THANK YOU", 64, WHITE, True, PP_ALIGN.CENTER, "Arial Black")
    add_textbox(sl, Inches(0), Inches(4.5), SLIDE_WIDTH, Inches(0.5),
                url, 18, RGBColor(0xBB, 0xDD, 0xBB), False, PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════
#  TEXT-ONLY PRESENTATION  (presentazione.pptx)
# ═══════════════════════════════════════════════════════

def build_text_pptx():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank = prs.slide_layouts[6]

    # ── 1: COVER ──
    build_cover(prs.slides.add_slide(blank))

    # ── 2: AGENDA ──
    sl = prs.slides.add_slide(blank)
    add_textbox(sl, Inches(0), Inches(0.8), SLIDE_WIDTH, Inches(0.7),
                "AGENDA", 30, BLACK, True, PP_ALIGN.CENTER, "Arial Black")
    agenda_items = [
        ("1", "La sfida: perché un LLM generico non basta"),
        ("2", "Setup: 4 agenti a confronto"),
        ("3", "Prompt 1 — Preparazione alla riunione"),
        ("4", "Prompt 2 — Riassunto email non lette"),
        ("5", "Prompt 3 — Azione: invio email"),
        ("6", "Riepilogo e messaggi chiave"),
    ]
    list_w = Inches(8)
    list_x = (SLIDE_WIDTH - list_w) // 2
    _, tf = add_rich_textbox(sl, list_x, Inches(2.2), list_w, Inches(4.5))
    tf.paragraphs[0].space_after = Pt(0)
    for num, text in agenda_items:
        p = add_paragraph(tf, font_size=18, space_after=Pt(18))
        add_run(p, f"{num}   ", 24, GREEN, True, font_name="Arial Black")
        add_run(p, text, 18, BLACK)

    # ── 3: LA SFIDA ──
    sl = prs.slides.add_slide(blank)
    build_section_slide(sl, "LA SFIDA",
                        "Un modello linguistico senza contesto aziendale è come un consulente brillante che non ha mai visto i tuoi dati.")

    # ── 4: STATEMENT ──
    build_statement(prs.slides.add_slide(blank))

    # ── 5: 4 LIVELLI ──
    build_four_pillars(prs.slides.add_slide(blank))

    # ── 6: PROMPT 1 ──
    sl = prs.slides.add_slide(blank)
    build_section_slide(sl, "PROMPT 1",
                        '"Cosa devo sapere prima della mia prossima riunione?"')

    # ── 7: P1 Agent 1 vs 2 ──
    sl = prs.slides.add_slide(blank)
    build_comparison_slide(sl,
        "PROMPT 1 — RISPOSTE A CONFRONTO", "BASE vs CONNESSIONI",
        ("🔴", "1", "AGENTE BASE", RED,
         "Non risultano informazioni specifiche sulla tua prossima riunione nei dati disponibili. "
         "Per offrirti un briefing completo, avrei bisogno di dettagli come l'orario, i partecipanti o l'argomento della riunione.",
         "🔴 Educato ma inutile. Non sa nulla.", RED),
        ("🟡", "2", "CONNESSIONI", YELLOW,
         "Horizon 3.0 - Stand up tecnico\n"
         "Oggi, 09:30 – 10:00 · Microsoft Teams\n"
         "Organizzatore: Marco Bianchi\n"
         "Partecipanti: Elena Rossi ✓, Luca Moretti, Sara Conti, Davide Rizzo (in attesa)\n"
         "Tu: accettato provvisoriamente",
         "🟡 Vede il calendario. Dati puntuali.", YELLOW))

    # ── 8: P1 Agent 3 vs 4 ──
    sl = prs.slides.add_slide(blank)
    build_comparison_slide(sl,
        "PROMPT 1 — RISPOSTE A CONFRONTO", "WORK IQ COPILOT vs COMPLETO",
        ("🟣", "3", "WORK IQ COPILOT", PURPLE,
         "📅 Meridian - Supply Chain Status - Weekly\n"
         "Domani 08:30 – 09:00 · Microsoft Teams\n"
         "Organizzatore: Paolo Gentile\n\n"
         "🧩 Meeting di stato settimanale su Meridian / Supply Chain Platform\n"
         "🔍 Parte di un filone di meeting (Phoenix Stand-Up, Horizon SAL, Atlas Weekly) — preparati su rischi e dipendenze",
         "🟣 Analisi strutturata, contesto storico, suggerimenti proattivi", PURPLE),
        ("🟢", "4", "COMPLETO", GREEN,
         "Horizon 3.0 - Stand up tecnico\n"
         "09:30 – 10:00 · Microsoft Teams\n"
         "Organizzatore: Marco Bianchi\n"
         "Elena Rossi ✓ · Luca, Sara, Davide (in attesa)\n"
         "Tu: tentativo\n\n"
         "💡 Suggerimento: prepara aggiornamenti rapidi su attività Horizon 3.0 ed eventuali blocchi.",
         "🟢 Dettagli operativi + stato RSVP + azioni concrete", GREEN))

    # ── 9: PROMPT 2 ──
    sl = prs.slides.add_slide(blank)
    build_section_slide(sl, "PROMPT 2",
                        '"Riassumi le email non lette più importanti di oggi."')

    # ── 10: P2 Agent 1 vs 2 ──
    sl = prs.slides.add_slide(blank)
    build_comparison_slide(sl,
        "PROMPT 2 — RISPOSTE A CONFRONTO", "BASE vs CONNESSIONI",
        ("🔴", "1", "AGENTE BASE", RED,
         "Non risultano email non lette importanti ricevute oggi nei dati disponibili. "
         "Se hai bisogno di un riepilogo su un periodo diverso o su una categoria specifica di email, fammi sapere!",
         "🔴 Nessun accesso alla mail", RED),
        ("🟡", "2", "CONNESSIONI", YELLOW,
         "1. PR: fix_validation_pipeline into main (Azure DevOps)\n"
         "2. Annullato: Sprint Review - Weekly Sync (Stefano Galli)\n"
         "3. [API Gateway] Integrazione Meridian — kick-off tecnico (Francesca Colombo)\n"
         "4. Analytics Dashboard [TEST] - Issue Report (errore critico)",
         "🟡 Elenca email categorizzate. Non distingue le priorità.", YELLOW))

    # ── 11: P2 Agent 3 vs 4 ──
    sl = prs.slides.add_slide(blank)
    build_comparison_slide(sl,
        "PROMPT 2 — RISPOSTE A CONFRONTO", "WORK IQ COPILOT vs COMPLETO",
        ("🟣", "3", "WORK IQ COPILOT", PURPLE,
         "Non risultano email non lette di oggi.\n\n"
         "Posso però:\n"
         "▪ Cercare nelle ultime 24-48 ore\n"
         "▪ Includere email già lette ma rilevanti\n"
         "▪ Filtrare per mittenti chiave (team Phoenix, Marco, Chiara)",
         "🟣 Spiega il contesto, propone strategie alternative intelligenti", PURPLE),
        ("🟢", "4", "COMPLETO", GREEN,
         "▪ Roberto Neri (29 mag) — Invito riunione Teams\n"
         "▪ Deskbooking (29 mag) — Check-in Reminder ×2\n"
         "▪ Microsoft Outlook (28 mag) — Non recapitabile: Follow-up con Sofia Marchetti ×2\n\n"
         "Vuoi che approfondisca per estrarre azioni?",
         "🟢 Email filtrate con data, anteprima e azioni di follow-up", GREEN))

    # ── 12: PROMPT 3 ──
    sl = prs.slides.add_slide(blank)
    build_section_slide(sl, "PROMPT 3",
                        '"Manda una email a s.marchetti@reply.it: digli che ho visto le sue domande sul progetto e ne parliamo alla prossima riunione."')

    # ── 13: PROMPT 3 — AZIONE ──
    sl = prs.slides.add_slide(blank)
    add_slide_title_center(sl, "AGENTE 4 — AZIONE COMPLETATA", "CAPIRE + AGIRE = VALORE REALE")

    card_w = Inches(7)
    card_x = (SLIDE_WIDTH - card_w) // 2
    card_top = Inches(2.0)
    card_h = Inches(3.8)
    add_rounded_rect(sl, card_x, card_top, card_w, card_h, LIGHT_GRAY, GREEN)
    pad = Inches(0.5)
    _, tf = add_rich_textbox(sl, card_x + pad, card_top + Inches(0.3), card_w - pad * 2, card_h - Inches(0.5))
    tf.paragraphs[0].space_after = Pt(0)
    add_paragraph(tf, "📨 A: s.marchetti@reply.it", 17, GREEN, True, space_after=Pt(12))
    add_paragraph(tf, 'Oggetto: "Discussione delle domande sul progetto"', 13, GRAY, italic=True, space_after=Pt(20))
    add_paragraph(tf, "Ciao Sofia,", 14, BLACK, space_after=Pt(6))
    add_paragraph(tf, "ho visto le tue domande sul progetto.", 14, BLACK, space_after=Pt(6))
    add_paragraph(tf, "Ne parleremo nella prossima riunione.", 14, BLACK, space_after=Pt(20))
    add_paragraph(tf, "✅ Email inviata con successo", 17, GREEN, True, space_after=Pt(0))

    add_textbox(sl, Inches(2), Inches(6.3), Inches(9.3), Inches(0.7),
                "L'agente ha compreso il contesto dalla conversazione, formulato un messaggio appropriato "
                "e inviato l'email — tutto in un unico passaggio.",
                13, GRAY, False, PP_ALIGN.CENTER)

    # ── 14: TABELLA ──
    build_table_slide(prs.slides.add_slide(blank))

    # ── 15: MESSAGGI CHIAVE ──
    build_messages_slide(prs.slides.add_slide(blank))

    # ── 16: THANK YOU ──
    build_thankyou(prs.slides.add_slide(blank))

    out_path = os.path.join(SCRIPT_DIR, "presentazione.pptx")
    prs.save(out_path)
    print(f"✅ Salvato: {out_path}")
    return out_path


# ═══════════════════════════════════════════════════════
#  VISUAL PRESENTATION  (presentazione_visual.pptx)
# ═══════════════════════════════════════════════════════

def build_visual_pptx():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank = prs.slide_layouts[6]

    def screenshot_slide(title, subtitle, img_file):
        """Full-width screenshot on dark background."""
        sl = prs.slides.add_slide(blank)
        set_slide_bg(sl, BG_DARK)
        add_textbox(sl, Inches(0), Inches(0.15), SLIDE_WIDTH, Inches(0.45),
                    title, 20, WHITE, True, PP_ALIGN.CENTER, "Arial Black")
        add_textbox(sl, Inches(0), Inches(0.6), SLIDE_WIDTH, Inches(0.3),
                    subtitle, 10, MUTED_TEXT, True, PP_ALIGN.CENTER, "Arial Black")
        img_path = os.path.join(SCREENSHOTS_DIR, img_file)
        if os.path.exists(img_path):
            add_image_fitted(sl, img_path, Inches(13.0), Inches(6.3), Inches(1.1))
        return sl

    def duo_slide(title, subtitle, img1, label1, badge1, img2, label2, badge2):
        """Two screenshots side by side on dark background."""
        sl = prs.slides.add_slide(blank)
        set_slide_bg(sl, BG_DARK)
        add_textbox(sl, Inches(0), Inches(0.15), SLIDE_WIDTH, Inches(0.45),
                    title, 18, WHITE, True, PP_ALIGN.CENTER, "Arial Black")
        add_textbox(sl, Inches(0), Inches(0.6), SLIDE_WIDTH, Inches(0.3),
                    subtitle, 10, MUTED_TEXT, True, PP_ALIGN.CENTER, "Arial Black")

        img_w = Inches(6.2)
        max_img_h = Inches(5.2)
        gap = Inches(0.3)
        total_w = img_w * 2 + gap
        left_x = (SLIDE_WIDTH - total_w) // 2
        right_x = left_x + img_w + gap
        top_y = Inches(1.1)

        for img_file, label, (badge_text, badge_color), x in [
            (img1, label1, badge1, left_x),
            (img2, label2, badge2, right_x),
        ]:
            img_path = os.path.join(SCREENSHOTS_DIR, img_file)
            actual_h = max_img_h  # fallback
            if os.path.exists(img_path):
                from PIL import Image as PILImage
                with PILImage.open(img_path) as pimg:
                    pw, ph = pimg.size
                aspect = pw / ph
                actual_w = img_w
                actual_h = int(actual_w / aspect)
                if actual_h > max_img_h:
                    actual_h = max_img_h
                    actual_w = int(actual_h * aspect)
                # Center image within column
                img_left = x + (img_w - actual_w) // 2
                sl.shapes.add_picture(img_path, img_left, top_y, actual_w, actual_h)

            # Badge — positioned right below image
            badge_y = top_y + actual_h + Inches(0.1)
            badge_rect = add_rounded_rect(sl, x, badge_y, Inches(2.2), Inches(0.3), badge_color)
            badge_rect.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            run = badge_rect.text_frame.paragraphs[0].add_run()
            run.text = badge_text
            run.font.size = Pt(8)
            run.font.color.rgb = WHITE
            run.font.bold = True
            run.font.name = "Arial Black"

            # Label — below badge
            add_textbox(sl, x, badge_y + Inches(0.35), img_w, Inches(0.4),
                        label, 10, RGBColor(0xBB, 0xBB, 0xBB))

        return sl

    # ── 1: COVER ──
    build_cover(prs.slides.add_slide(blank))

    # ── 2: STATEMENT ──
    build_statement(prs.slides.add_slide(blank))

    # ── 3: 4 LIVELLI ──
    build_four_pillars(prs.slides.add_slide(blank))

    # ── 4: DASHBOARD ──
    screenshot_slide("DASHBOARD DI PARTENZA",
                     "4 AGENTI COPILOT STUDIO CONNESSI IN PARALLELO",
                     "00_dashboard_pulita.png")

    # ── 5: PROMPT 1 SECTION ──
    sl = prs.slides.add_slide(blank)
    build_section_slide(sl, "PROMPT 1",
                        '"Cosa devo sapere prima della mia prossima riunione?"')

    # ── 6: P1 OVERVIEW ──
    screenshot_slide("PROMPT 1 — VISTA COMPARATIVA",
                     "STESSO PROMPT, 4 RISPOSTE DIVERSE",
                     "01_prompt1_overview.png")

    # ── 7: P1 AGENT 1 vs 2 ──
    duo_slide("PROMPT 1 — BASE vs CONNESSIONI",
              'DA "NON HO INFORMAZIONI" A DETTAGLI COMPLETI',
              "02_prompt1_agent1.png", "🔴 Educato ma inutile. Non sa nulla del tuo lavoro.",
              ("1  BASE", RED),
              "02_prompt1_agent2.png", "🟡 Vede il calendario: titolo, orario, partecipanti.",
              ("2  CONNESSIONI", YELLOW))

    # ── 8: P1 AGENT 3 vs 4 ──
    duo_slide("PROMPT 1 — COPILOT vs COMPLETO",
              "ANALISI PROFONDA E SUGGERIMENTI OPERATIVI",
              "02_prompt1_agent3.png", "🟣 Analisi strutturata, contesto storico, briefing operativo.",
              ("3  WORK IQ COPILOT", PURPLE),
              "02_prompt1_agent4.png", "🟢 Dettagli operativi + stato RSVP + azioni concrete.",
              ("4  COMPLETO", GREEN))

    # ── 9: PROMPT 2 SECTION ──
    sl = prs.slides.add_slide(blank)
    build_section_slide(sl, "PROMPT 2",
                        '"Riassumi le email non lette più importanti di oggi."')

    # ── 10: P2 OVERVIEW ──
    screenshot_slide("PROMPT 2 — VISTA COMPARATIVA",
                     "EMAIL: DA ZERO INFORMAZIONI A INSIGHT STRUTTURATI",
                     "06_prompt2_overview.png")

    # ── 11: P2 AGENT 1 vs 2 ──
    duo_slide("PROMPT 2 — BASE vs CONNESSIONI",
              "ACCESSO AI DATI VS ANALISI CIECA",
              "07_prompt2_agent1.png", "🔴 Nessun accesso alla mail.",
              ("1  BASE", RED),
              "07_prompt2_agent2.png", "🟡 Elenca email categorizzate con azioni suggerite.",
              ("2  CONNESSIONI", YELLOW))

    # ── 12: P2 AGENT 3 vs 4 ──
    duo_slide("PROMPT 2 — COPILOT vs COMPLETO",
              "STRATEGIE ALTERNATIVE VS FILTRAGGIO INTELLIGENTE",
              "07_prompt2_agent3.png", "🟣 Propone strategie alternative intelligenti.",
              ("3  WORK IQ COPILOT", PURPLE),
              "07_prompt2_agent4.png", "🟢 Email filtrate per rilevanza con data e anteprima.",
              ("4  COMPLETO", GREEN))

    # ── 13: PROMPT 3 SECTION ──
    sl = prs.slides.add_slide(blank)
    build_section_slide(sl, "PROMPT 3",
                        '"Manda una email a s.marchetti@reply.it: digli che ho visto le sue domande sul progetto e ne parliamo alla prossima riunione."')

    # ── 14: P3 AGENT 4 — Annotated ──
    sl = prs.slides.add_slide(blank)
    set_slide_bg(sl, BG_DARK)
    add_textbox(sl, Inches(0), Inches(0.15), SLIDE_WIDTH, Inches(0.45),
                "AGENTE 4 — AZIONE COMPLETATA", 20, WHITE, True, PP_ALIGN.CENTER, "Arial Black")
    add_textbox(sl, Inches(0), Inches(0.6), SLIDE_WIDTH, Inches(0.3),
                "CAPIRE + AGIRE = VALORE REALE", 10, MUTED_TEXT, True, PP_ALIGN.CENTER, "Arial Black")

    # Screenshot (left 2/3)
    img_path = os.path.join(SCREENSHOTS_DIR, "12_prompt3_agent4_email.png")
    if os.path.exists(img_path):
        from PIL import Image as PILImage
        with PILImage.open(img_path) as pimg:
            pw, ph = pimg.size
        aspect = pw / ph
        img_display_w = Inches(9.0)
        img_display_h = int(img_display_w / aspect)
        max_h = Inches(6.2)
        if img_display_h > max_h:
            img_display_h = max_h
            img_display_w = int(img_display_h * aspect)
        sl.shapes.add_picture(img_path, Inches(0.4), Inches(1.1), img_display_w, img_display_h)

    # Annotation (right 1/3)
    ann_x = Inches(9.8)
    ann_w = Inches(3.2)
    _, tf = add_rich_textbox(sl, ann_x, Inches(1.8), ann_w, Inches(5))
    tf.paragraphs[0].space_after = Pt(0)
    add_paragraph(tf, "UN SINGOLO PROMPT.", 17, GREEN, True, font_name="Arial Black", space_after=Pt(16))
    add_paragraph(tf, "L'agente ha compreso il contesto dalla conversazione, formulato un messaggio appropriato e inviato l'email.",
                  13, MUTED_TEXT, space_after=Pt(28))
    p = add_paragraph(tf, "", space_after=Pt(4))
    add_run(p, "1", 52, WHITE, True, font_name="Arial Black")
    add_paragraph(tf, "PASSAGGIO PER\nCAPIRE E AGIRE", 10, MUTED_TEXT, True, font_name="Arial Black", space_after=Pt(28))
    add_paragraph(tf, "📨  A: s.marchetti@reply.it", 12, WHITE, space_after=Pt(8))
    add_paragraph(tf, "✅  Email inviata con successo", 14, GREEN, True, space_after=Pt(0))

    # ── 15: PANORAMICA ──
    screenshot_slide("PANORAMICA FINALE",
                     "TUTTA LA CONVERSAZIONE A CONFRONTO",
                     "13_panoramica_finale.png")

    # ── 16: TABELLA ──
    build_table_slide(prs.slides.add_slide(blank))

    # ── 17: MESSAGGI CHIAVE ──
    build_messages_slide(prs.slides.add_slide(blank))

    # ── 18: THANK YOU ──
    build_thankyou(prs.slides.add_slide(blank), "www.reply.com")

    out_path = os.path.join(SCRIPT_DIR, "presentazione_visual.pptx")
    prs.save(out_path)
    print(f"✅ Salvato: {out_path}")
    return out_path


# ═══════════════════════════════════════════════════════
if __name__ == "__main__":
    print("Generazione PPTX in corso...\n")
    build_text_pptx()
    build_visual_pptx()
    print("\n🎉 Entrambi i file PPTX sono stati generati!")
